import base64
import codecs
import csv
import fnmatch
import html as html_lib
import io
import json
import logging
import mimetypes
import os
import re
import shutil
import tarfile
import time
import uuid
import zipfile
from datetime import datetime, timedelta
import ipaddress
import socket
from html.parser import HTMLParser
import redis
import urllib.error
import urllib.parse
import urllib.request
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import requests

from helpers import redis_blob_client, redis_client
from verba_loader import load_verbas_from_directory
from verba_registry import reload_verbas
from verba_result import action_failure, action_success
from notify import dispatch_notification_sync, notifier_destination_catalog
from notify.queue import ALLOWED_PLATFORMS, normalize_platform as normalize_notify_platform
from vision_settings import (
    DEFAULT_VISION_API_BASE,
    DEFAULT_VISION_MODEL,
    get_vision_settings,
)


BASE_DIR = Path(__file__).resolve().parent
_agent_root_env = str(os.getenv("TATER_AGENT_ROOT", "") or "").strip()
if _agent_root_env:
    _agent_root_path = Path(_agent_root_env).expanduser()
    if not _agent_root_path.is_absolute():
        _agent_root_path = BASE_DIR / _agent_root_path
else:
    _agent_root_path = BASE_DIR / "agent_lab"
AGENT_LAB_DIR = _agent_root_path.resolve()
AGENT_PLUGINS_DIR = AGENT_LAB_DIR / "verba"
AGENT_PORTALS_DIR = AGENT_LAB_DIR / "portals"
AGENT_ARTIFACTS_DIR = AGENT_LAB_DIR / "artifacts"
AGENT_DOCUMENTS_DIR = AGENT_LAB_DIR / "documents"
AGENT_DOWNLOADS_DIR = AGENT_LAB_DIR / "downloads"
AGENT_WORKSPACE_DIR = AGENT_LAB_DIR / "workspace"
AGENT_LOGS_DIR = AGENT_LAB_DIR / "logs"
AGENT_REQUIREMENTS = AGENT_LAB_DIR / "requirements.txt"

STABLE_PLUGINS_DIR = BASE_DIR / os.getenv("TATER_VERBA_DIR", "verba")
STABLE_PORTALS_DIR = BASE_DIR / "portals"

_SAFE_NAME_RE = re.compile(r"^[a-zA-Z0-9_\-]+$")
_READ_FILE_MAX_CHARS = int(os.getenv("TATER_READ_FILE_MAX_CHARS", "400000"))
_READ_PDF_MAX_PAGES = int(os.getenv("TATER_READ_PDF_MAX_PAGES", "120"))
_READ_DOCX_MAX_TABLE_ROWS = int(os.getenv("TATER_READ_DOCX_MAX_TABLE_ROWS", "300"))
_READ_XLSX_MAX_SHEETS = int(os.getenv("TATER_READ_XLSX_MAX_SHEETS", "4"))
_READ_XLSX_MAX_ROWS_PER_SHEET = int(os.getenv("TATER_READ_XLSX_MAX_ROWS_PER_SHEET", "120"))
_READ_XLSX_MAX_COLS_PER_ROW = int(os.getenv("TATER_READ_XLSX_MAX_COLS_PER_ROW", "24"))
_READ_CSV_MAX_ROWS = int(os.getenv("TATER_READ_CSV_MAX_ROWS", "500"))
_READ_CSV_MAX_COLS = int(os.getenv("TATER_READ_CSV_MAX_COLS", "64"))
_READ_PPTX_MAX_SLIDES = int(os.getenv("TATER_READ_PPTX_MAX_SLIDES", "120"))
_SEARCH_DEFAULT_MAX_RESULTS = int(os.getenv("TATER_SEARCH_MAX_RESULTS", "100"))
_SEARCH_MAX_FILE_CHARS = int(os.getenv("TATER_SEARCH_MAX_FILE_CHARS", "200000"))
_ARCHIVE_LIST_MAX_ENTRIES = int(os.getenv("TATER_ARCHIVE_LIST_MAX_ENTRIES", "1000"))
_ARCHIVE_EXTRACT_MAX_FILES = int(os.getenv("TATER_ARCHIVE_EXTRACT_MAX_FILES", "1000"))

WEB_SEARCH_API_KEY_REDIS_KEY = "tater:web_search:google_api_key"
WEB_SEARCH_CX_REDIS_KEY = "tater:web_search:google_cx"
WEB_SEARCH_LEGACY_SETTINGS_KEY = "verba_settings:Web Search"
WEB_SEARCH_TIMEOUT_SEC = int(os.getenv("TATER_WEB_SEARCH_TIMEOUT_SEC", "15"))
WEB_SEARCH_MAX_RESULTS = int(os.getenv("TATER_WEB_SEARCH_MAX_RESULTS", "10"))
WEB_SEARCH_MAX_RESPONSE_BYTES = int(os.getenv("TATER_WEB_SEARCH_MAX_RESPONSE_BYTES", "2000000"))
WEB_SEARCH_MAX_SNIPPET_CHARS = int(os.getenv("TATER_WEB_SEARCH_MAX_SNIPPET_CHARS", "600"))
WEBPAGE_TIMEOUT_SEC = int(os.getenv("TATER_WEBPAGE_TIMEOUT_SEC", "20"))
WEBPAGE_MAX_RESPONSE_BYTES = int(os.getenv("TATER_WEBPAGE_MAX_RESPONSE_BYTES", "4000000"))
WEBPAGE_MAX_TEXT_CHARS = int(os.getenv("TATER_WEBPAGE_MAX_TEXT_CHARS", "180000"))
WEBPAGE_MAX_PREVIEW_CHARS = int(os.getenv("TATER_WEBPAGE_MAX_PREVIEW_CHARS", "1400"))
WEBPAGE_BLOCK_STATUS_CODES = {403, 406, 409, 412, 418, 425, 429, 451, 503}
WEBPAGE_BOT_BLOCK_MARKERS = (
    "access denied",
    "attention required",
    "verify you are human",
    "verification required",
    "captcha",
    "cloudflare",
    "request blocked",
    "bot detection",
    "checking your browser",
)
WEBPAGE_USER_AGENTS = (
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) "
    "Chrome/125.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 14_5) AppleWebKit/605.1.15 (KHTML, like Gecko) "
    "Version/17.5 Safari/605.1.15",
    "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) "
    "Chrome/124.0.0.0 Safari/537.36",
)
VISION_ALLOWED_MIMETYPES = {
    "image/png",
    "image/jpeg",
    "image/jpg",
    "image/webp",
    "image/gif",
    "image/bmp",
    "image/tiff",
}
VISION_DEFAULT_PROMPT = (
    "Describe this image clearly and concisely. Mention important objects, people, actions, "
    "and any visible text."
)
WEBUI_FILE_BLOB_KEY_PREFIX = "webui:file:"

AI_TASKS_KEY_PREFIX = "reminders:"
AI_TASKS_DUE_ZSET = "reminders:due"
AI_TASKS_DAILY_MARKERS = ("every day", "everyday", "daily", "each day")
AI_TASKS_WEEKLY_MARKERS = ("every week", "each week", "weekly")
AI_TASKS_MONTHLY_MARKERS = ("every month", "each month", "monthly")
AI_TASKS_WEEKDAY_MAP = {
    "monday": 0,
    "mon": 0,
    "mondays": 0,
    "tuesday": 1,
    "tue": 1,
    "tues": 1,
    "tuesdays": 1,
    "wednesday": 2,
    "wed": 2,
    "weds": 2,
    "wedsday": 2,
    "wednesdays": 2,
    "thursday": 3,
    "thu": 3,
    "thur": 3,
    "thurs": 3,
    "thursdays": 3,
    "friday": 4,
    "fri": 4,
    "fridays": 4,
    "saturday": 5,
    "sat": 5,
    "saturdays": 5,
    "sunday": 6,
    "sun": 6,
    "sundays": 6,
}
AI_TASKS_LOCAL_TZ_HINT_RE = re.compile(r"\bassume\s+local\s+timezone\b\.?", re.IGNORECASE)
AI_TASKS_WEATHER_DEFAULT_HINT_RE = re.compile(
    r"\bif\s+location\s+is\s+not\s+specified,\s+use\s+the\s+configured\s+default\s+weather\s+location\b\.?",
    re.IGNORECASE,
)
AI_TASKS_TIME_PREFIX_RE = re.compile(
    r"^\s*(?:(?:in|after)\s+\d+\s*(?:seconds?|minutes?|hours?|days?|weeks?)|at\s+(?:\d{3,4}|\d{1,2}(?::\d{2})?(?::\d{2})?)\s*(?:am|pm)?)\b",
    re.IGNORECASE,
)
AI_TASKS_SCHEDULE_PREFIX_PATTERNS = (
    re.compile(
        r"^\s*(?:every\s+day|everyday|daily|each\s+day|weekdays?|weekends?)\b(?:\s+at\s+(?:\d{3,4}|\d{1,2}(?::\d{2})?(?::\d{2})?)\s*(?:am|pm)?)?\s*(?:,|:|-)?\s*",
        re.IGNORECASE,
    ),
    re.compile(
        r"^\s*(?:every\s+week|each\s+week|weekly)\b(?:\s+on\s+[a-z,\s]+)?(?:\s+at\s+(?:\d{3,4}|\d{1,2}(?::\d{2})?(?::\d{2})?)\s*(?:am|pm)?)?\s*(?:,|:|-)?\s*",
        re.IGNORECASE,
    ),
    re.compile(
        r"^\s*on\s+(?:mon(?:day|days?)?|tues?(?:day|days?)?|wed(?:nesday|nesdays|s|sday|sdays)?|thu(?:r|rs|rsday|rsdays|day|days)?|fri(?:day|days)?|sat(?:urday|urdays)?|sun(?:day|days)?)(?:\s*(?:,|and)\s*(?:mon(?:day|days?)?|tues?(?:day|days?)?|wed(?:nesday|nesdays|s|sday|sdays)?|thu(?:r|rs|rsday|rsdays|day|days)?|fri(?:day|days)?|sat(?:urday|urdays)?|sun(?:day|days)?))*\s+(?:every|each)\s+week\b(?:\s+at\s+(?:\d{3,4}|\d{1,2}(?::\d{2})?(?::\d{2})?)\s*(?:am|pm)?)?\s*(?:,|:|-)?\s*",
        re.IGNORECASE,
    ),
    re.compile(
        r"^\s*on\s+(?:the\s+)?(?:[12]?\d|3[01])(?:st|nd|rd|th)(?:\s*(?:,|and)\s*(?:the\s+)?(?:[12]?\d|3[01])(?:st|nd|rd|th))*\s+of\s+(?:every|each)\s+month\b(?:\s+at\s+(?:\d{3,4}|\d{1,2}(?::\d{2})?(?::\d{2})?)\s*(?:am|pm)?)?\s*(?:,|:|-)?\s*",
        re.IGNORECASE,
    ),
    re.compile(
        r"^\s*(?:every\s+month|each\s+month|monthly)\b(?:\s+on\s+(?:the\s+)?(?:[12]?\d|3[01])(?:st|nd|rd|th)(?:\s*(?:,|and)\s*(?:the\s+)?(?:[12]?\d|3[01])(?:st|nd|rd|th))*)?(?:\s+at\s+(?:\d{3,4}|\d{1,2}(?::\d{2})?(?::\d{2})?)\s*(?:am|pm)?)?\s*(?:,|:|-)?\s*",
        re.IGNORECASE,
    ),
    re.compile(
        r"^\s*every\s+\d+\s*(?:seconds?|minutes?|hours?|days?|weeks?)\b\s*(?:,|:|-)?\s*",
        re.IGNORECASE,
    ),
    re.compile(
        r"^\s*(?:in|after)\s+\d+\s*(?:seconds?|minutes?|hours?|days?|weeks?)\b\s*(?:,|:|-)?\s*",
        re.IGNORECASE,
    ),
    re.compile(
        r"^\s*at\s+(?:\d{3,4}|\d{1,2}(?::\d{2})?(?::\d{2})?)\s*(?:am|pm)?\b\s*(?:,|:|-)?\s*",
        re.IGNORECASE,
    ),
)


logger = logging.getLogger("kernel_tools")


def _ensure_dirs() -> None:
    for path in (
        AGENT_LAB_DIR,
        AGENT_PLUGINS_DIR,
        AGENT_PORTALS_DIR,
        AGENT_ARTIFACTS_DIR,
        AGENT_DOCUMENTS_DIR,
        AGENT_DOWNLOADS_DIR,
        AGENT_WORKSPACE_DIR,
        AGENT_LOGS_DIR,
    ):
        path.mkdir(parents=True, exist_ok=True)
    if not AGENT_REQUIREMENTS.exists():
        AGENT_REQUIREMENTS.write_text("", encoding="utf-8")


def _log_write(action: str, path: Path, size: int = 0) -> None:
    try:
        _ensure_dirs()
        ts = time.strftime("%Y-%m-%d %H:%M:%S")
        line = f"{ts} | {action} | {path} | {size} bytes\n"
        with (AGENT_LOGS_DIR / "agent_writes.log").open("a", encoding="utf-8") as f:
            f.write(line)
    except Exception:
        return


def _resolve_safe_path(path: str, allowed_roots: List[Path]) -> Optional[Path]:
    if not path:
        return None

    raw = str(path).strip()
    normalized = raw.replace("\\", "/")
    while normalized.startswith("./"):
        normalized = normalized[2:]

    # Virtual workspace root aliases (agent sees root as /).
    if normalized in {"download", "downloads"}:
        normalized = "downloads"
    elif normalized.startswith("download/"):
        normalized = "downloads/" + normalized[len("download/") :]
    elif normalized in {"document", "documents"}:
        normalized = "documents"
    elif normalized.startswith("document/"):
        normalized = "documents/" + normalized[len("document/") :]
    elif normalized in {"/download", "/downloads"}:
        normalized = "/downloads"
    elif normalized.startswith("/download/"):
        normalized = "/downloads/" + normalized[len("/download/") :]
    elif normalized in {"/document", "/documents"}:
        normalized = "/documents"
    elif normalized.startswith("/document/"):
        normalized = "/documents/" + normalized[len("/document/") :]

    if normalized in {"/", "/."}:
        raw = str(AGENT_LAB_DIR)
    elif normalized == "/agent_lab":
        raw = str(AGENT_LAB_DIR)
    elif normalized.startswith("/agent_lab/"):
        suffix = normalized[len("/agent_lab/") :]
        raw = str(AGENT_LAB_DIR / suffix)
    elif normalized.startswith("/"):
        raw = str(AGENT_LAB_DIR / normalized.lstrip("/"))
    else:
        raw = str(AGENT_LAB_DIR / normalized)

    p = Path(raw)
    if not p.is_absolute():
        p = (BASE_DIR / p).resolve()
    else:
        p = p.resolve()

    for root in allowed_roots:
        try:
            root_resolved = root.resolve()
        except Exception:
            root_resolved = root
        if p == root_resolved or root_resolved in p.parents:
            return p
    return None


def _display_workspace_path(path: Any) -> str:
    raw = str(path or "").strip()
    if not raw:
        return ""
    try:
        root = AGENT_LAB_DIR.resolve()
        resolved = Path(raw).resolve()
        if resolved == root:
            return "/"
        if root in resolved.parents:
            rel = resolved.relative_to(root).as_posix()
            return f"/{rel}" if rel else "/"
    except Exception:
        pass
    return raw


def _coerce_int(value: Any, default: int, min_value: int = 0, max_value: Optional[int] = None) -> int:
    try:
        out = int(float(value))
    except Exception:
        out = int(default)
    if out < min_value:
        out = min_value
    if max_value is not None and out > max_value:
        out = max_value
    return out


def _coerce_bool(value: Any, default: bool = False) -> bool:
    if value is None:
        return bool(default)
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return value != 0
    if isinstance(value, str):
        v = value.strip().lower()
        if v in {"1", "true", "yes", "on"}:
            return True
        if v in {"0", "false", "no", "off"}:
            return False
    return bool(default)


def _slice_content(text: str, start: Any = 0, max_chars: Any = None) -> Tuple[str, Dict[str, Any]]:
    start_i = _coerce_int(start, default=0, min_value=0)
    limit_default = _READ_FILE_MAX_CHARS
    limit_i = _coerce_int(
        max_chars if max_chars is not None else limit_default,
        default=limit_default,
        min_value=1,
        max_value=2_000_000,
    )
    total = len(text)
    if start_i > total:
        start_i = total
    end_i = min(total, start_i + limit_i)
    chunk = text[start_i:end_i]
    has_more = end_i < total
    meta = {
        "start": start_i,
        "end": end_i,
        "max_chars": limit_i,
        "total_chars": total,
        "returned_chars": len(chunk),
        "has_more": has_more,
        "next_start": end_i if has_more else None,
    }
    return chunk, meta


def _read_text(path: Path) -> str:
    raw = path.read_bytes()
    if not raw:
        return ""

    if raw.startswith(codecs.BOM_UTF8):
        text = raw.decode("utf-8-sig", errors="replace")
    elif raw.startswith(codecs.BOM_UTF16_LE) or raw.startswith(codecs.BOM_UTF16_BE):
        text = raw.decode("utf-16", errors="replace")
    elif raw.startswith(codecs.BOM_UTF32_LE) or raw.startswith(codecs.BOM_UTF32_BE):
        text = raw.decode("utf-32", errors="replace")
    else:
        if b"\x00" in raw[:8192]:
            raise ValueError(
                "Binary file is not readable as plain text. Supported document formats: "
                ".pdf, .docx, .xlsx, .xlsm, .csv, .tsv, .pptx."
            )
        text = raw.decode("utf-8", errors="replace")

    return text


def _read_pdf_text(path: Path) -> Tuple[str, Dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except Exception as e:
        raise RuntimeError("PDF parsing requires the `pypdf` package.") from e

    reader = PdfReader(str(path))
    total_pages = len(reader.pages)
    max_pages = max(1, _READ_PDF_MAX_PAGES)
    pages_to_read = min(total_pages, max_pages)
    chunks: List[str] = []

    for idx in range(pages_to_read):
        page_text = ""
        try:
            page_text = reader.pages[idx].extract_text() or ""
        except Exception:
            page_text = ""
        page_text = page_text.strip()
        if page_text:
            chunks.append(f"[Page {idx + 1}]\n{page_text}")
        else:
            chunks.append(f"[Page {idx + 1}]\n")

    if not chunks:
        chunks = ["[No extractable text found in PDF.]"]

    merged = "\n\n".join(chunks).strip()
    if not merged:
        merged = "[No extractable text found in PDF.]"
    metadata = {
        "format": "pdf",
        "pages": total_pages,
        "pages_read": pages_to_read,
        "source_truncated": bool(total_pages > pages_to_read),
    }
    return merged, metadata


def _read_docx_text(path: Path) -> Tuple[str, Dict[str, Any]]:
    try:
        import docx
    except Exception as e:
        raise RuntimeError("DOCX parsing requires the `python-docx` package.") from e

    doc = docx.Document(str(path))
    chunks: List[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            chunks.append(text)

    table_rows = 0
    total_table_rows = 0
    max_rows = max(1, _READ_DOCX_MAX_TABLE_ROWS)
    for table in doc.tables:
        for row in table.rows:
            total_table_rows += 1
            if table_rows >= max_rows:
                break
            cells = [str((cell.text or "")).strip() for cell in row.cells]
            if any(cells):
                chunks.append(" | ".join(cells))
            table_rows += 1

    merged = "\n".join(chunks).strip() or "[No extractable text found in DOCX.]"
    metadata = {
        "format": "docx",
        "paragraphs": len(doc.paragraphs),
        "table_rows": total_table_rows,
        "table_rows_read": min(table_rows, max_rows),
        "source_truncated": bool(total_table_rows > max_rows),
    }
    return merged, metadata


def _read_xlsx_text(path: Path) -> Tuple[str, Dict[str, Any]]:
    try:
        from openpyxl import load_workbook
    except Exception as e:
        raise RuntimeError("XLSX parsing requires the `openpyxl` package.") from e

    max_sheets = max(1, _READ_XLSX_MAX_SHEETS)
    max_rows = max(1, _READ_XLSX_MAX_ROWS_PER_SHEET)
    max_cols = max(1, _READ_XLSX_MAX_COLS_PER_ROW)
    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    chunks: List[str] = []
    sheet_names = list(wb.sheetnames or [])
    sheets_to_read = sheet_names[:max_sheets]
    source_truncated = False

    for sheet_name in sheets_to_read:
        ws = wb[sheet_name]
        chunks.append(f"[Sheet: {sheet_name}]")
        row_index = 0
        for row in ws.iter_rows(min_row=1, max_row=max_rows + 1, max_col=max_cols, values_only=True):
            if row_index >= max_rows:
                source_truncated = True
                break
            values = ["" if cell is None else str(cell) for cell in row]
            if any(v.strip() for v in values):
                chunks.append("\t".join(values).rstrip())
            row_index += 1
        if row_index == 0:
            chunks.append("[Empty sheet]")
        chunks.append("")

    wb.close()
    merged = "\n".join(chunks).strip() or "[No extractable content found in XLSX.]"
    source_truncated = bool(source_truncated or len(sheet_names) > len(sheets_to_read))
    metadata = {
        "format": "xlsx",
        "sheets": len(sheet_names),
        "sheets_read": len(sheets_to_read),
        "rows_per_sheet_limit": max_rows,
        "source_truncated": source_truncated,
    }
    return merged, metadata


def _sniff_csv_delimiter(sample: str, path: Path) -> str:
    if path.suffix.lower() == ".tsv":
        return "\t"
    try:
        dialect = csv.Sniffer().sniff(sample or "", delimiters=",;\t|")
        delim = getattr(dialect, "delimiter", ",")
        return delim if isinstance(delim, str) and delim else ","
    except Exception:
        return ","


def _read_csv_text(path: Path) -> Tuple[str, Dict[str, Any]]:
    try:
        raw = path.read_bytes()
    except Exception as e:
        raise RuntimeError(f"Unable to read CSV file: {e}") from e
    if not raw:
        return "", {"format": "csv", "rows": 0, "rows_read": 0, "source_truncated": False}

    try:
        text = raw.decode("utf-8-sig")
    except Exception:
        text = raw.decode("utf-8", errors="replace")

    delimiter = _sniff_csv_delimiter(text[:4096], path)
    max_rows = max(1, _READ_CSV_MAX_ROWS)
    max_cols = max(1, _READ_CSV_MAX_COLS)
    reader = csv.reader(io.StringIO(text, newline=""), delimiter=delimiter)

    out_lines: List[str] = []
    total_rows = 0
    rows_read = 0
    max_seen_cols = 0
    source_truncated = False

    for row in reader:
        total_rows += 1
        max_seen_cols = max(max_seen_cols, len(row))
        if rows_read >= max_rows:
            source_truncated = True
            continue
        if len(row) > max_cols:
            source_truncated = True
        visible = row[:max_cols]
        out_lines.append("\t".join(str(cell) for cell in visible))
        rows_read += 1

    merged = "\n".join(out_lines).strip()
    if not merged and total_rows == 0:
        merged = ""
    elif not merged:
        merged = "[No non-empty rows found in CSV/TSV.]"

    metadata = {
        "format": "csv",
        "delimiter": delimiter,
        "rows": total_rows,
        "rows_read": rows_read,
        "max_columns_seen": max_seen_cols,
        "source_truncated": bool(source_truncated),
    }
    return merged, metadata


def _read_pptx_text(path: Path) -> Tuple[str, Dict[str, Any]]:
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("PPTX parsing requires the `python-pptx` package.") from e

    prs = Presentation(str(path))
    total_slides = len(prs.slides)
    max_slides = max(1, _READ_PPTX_MAX_SLIDES)
    slides_to_read = min(total_slides, max_slides)
    chunks: List[str] = []

    for idx, slide in enumerate(prs.slides):
        if idx >= slides_to_read:
            break
        slide_lines: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                raw = shape.text or ""
                text = raw.strip()
                if text:
                    slide_lines.append(text)
            elif hasattr(shape, "table") and shape.table is not None:
                for row in shape.table.rows:
                    cells = [str((cell.text or "")).strip() for cell in row.cells]
                    if any(cells):
                        slide_lines.append(" | ".join(cells))

        if slide_lines:
            chunks.append(f"[Slide {idx + 1}]\n" + "\n".join(slide_lines))
        else:
            chunks.append(f"[Slide {idx + 1}]\n")

    merged = "\n\n".join(chunks).strip() or "[No extractable text found in PPTX.]"
    metadata = {
        "format": "pptx",
        "slides": total_slides,
        "slides_read": slides_to_read,
        "source_truncated": bool(total_slides > slides_to_read),
    }
    return merged, metadata


def _extract_file_content(path: Path) -> Tuple[str, Dict[str, Any]]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _read_pdf_text(path)
    if ext == ".docx":
        return _read_docx_text(path)
    if ext in {".xlsx", ".xlsm"}:
        return _read_xlsx_text(path)
    if ext in {".csv", ".tsv"}:
        return _read_csv_text(path)
    if ext == ".pptx":
        return _read_pptx_text(path)
    return _read_text(path), {"format": "text", "source_truncated": False}


def _sanitize_filename(name: str) -> str:
    raw = os.path.basename((name or "").strip())
    if not raw:
        return ""
    safe = re.sub(r"[^A-Za-z0-9._-]+", "_", raw)
    safe = safe.lstrip(".")
    return safe or ""


def _is_private_ip(ip_str: str) -> bool:
    try:
        ip = ipaddress.ip_address(ip_str)
    except ValueError:
        return False
    return (
        ip.is_private
        or ip.is_loopback
        or ip.is_link_local
        or ip.is_reserved
        or ip.is_multicast
    )


def _host_is_private(host: str) -> Tuple[bool, Optional[str]]:
    if not host:
        return True, "URL must include a host."
    try:
        if _is_private_ip(host):
            return True, None
    except Exception:
        return True, "Invalid host."
    try:
        infos = socket.getaddrinfo(host, None)
    except Exception:
        return True, "Unable to resolve host."
    for info in infos:
        ip_str = info[4][0]
        if _is_private_ip(ip_str):
            return True, None
    return False, None


def _validate_url(url: str) -> Optional[str]:
    if not url:
        return "URL is required."
    parsed = urllib.parse.urlparse(url)
    if parsed.scheme not in {"http", "https"}:
        return "Only http/https URLs are allowed."
    if not parsed.hostname:
        return "URL must include a host."
    is_private, err = _host_is_private(parsed.hostname)
    if err:
        return err
    if is_private:
        return "Private or local network hosts are not allowed."
    return None


def _normalize_url_input(url: Any) -> str:
    raw = _as_text(url).strip()
    if not raw:
        return ""
    if raw.startswith("//"):
        return f"https:{raw}"
    parsed = urllib.parse.urlparse(raw)
    if parsed.scheme:
        return raw
    if re.match(r"^(?:www\.)?[A-Za-z0-9][A-Za-z0-9.-]+\.[A-Za-z]{2,}(?:[/:?#].*)?$", raw):
        return f"https://{raw}"
    return raw


def _webpage_request_headers(user_agent: str) -> Dict[str, str]:
    ua = str(user_agent or "").strip() or WEBPAGE_USER_AGENTS[0]
    return {
        "User-Agent": ua,
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,text/plain;q=0.8,*/*;q=0.5",
        "Accept-Language": "en-US,en;q=0.9",
        "Accept-Encoding": "gzip, deflate, br",
        "Cache-Control": "no-cache",
        "Pragma": "no-cache",
        "DNT": "1",
        "Upgrade-Insecure-Requests": "1",
    }


def _webpage_is_textual_content_type(content_type: str) -> bool:
    token = str(content_type or "").split(";", 1)[0].strip().lower()
    if not token:
        return True
    return (
        token.startswith("text/")
        or token in {"application/xhtml+xml", "application/json", "application/xml"}
        or token.endswith("+json")
        or token.endswith("+xml")
    )


def _webpage_decode_text(raw: bytes, *, content_type: str, response_encoding: str = "") -> str:
    encoding = ""
    ct = str(content_type or "")
    match = re.search(r"charset=([A-Za-z0-9._:-]+)", ct, flags=re.IGNORECASE)
    if match:
        encoding = str(match.group(1) or "").strip().strip(";")
    if not encoding:
        encoding = str(response_encoding or "").strip()
    if not encoding:
        encoding = "utf-8"
    try:
        return raw.decode(encoding, errors="replace")
    except Exception:
        try:
            return raw.decode("utf-8", errors="replace")
        except Exception:
            return raw.decode("latin-1", errors="replace")


def _webpage_text_from_html(html: str, *, max_chars: int) -> str:
    text = str(html or "")
    if not text:
        return ""

    # Remove heavy non-content regions before stripping tags.
    text = re.sub(r"(?is)<(script|style|noscript|template|svg).*?>.*?</\1>", " ", text)
    text = re.sub(r"(?is)<br\s*/?>", "\n", text)
    text = re.sub(r"(?is)</(p|div|section|article|aside|main|li|tr|td|th|h1|h2|h3|h4|h5|h6)>", "\n", text)
    text = re.sub(r"(?is)<[^>]+>", " ", text)
    text = html_lib.unescape(text)
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t\f\v]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = "\n".join(part.strip() for part in text.split("\n"))
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return ""
    if len(text) > max_chars:
        clipped = text[:max_chars]
        if " " in clipped[1200:]:
            clipped = clipped[: clipped.rfind(" ")]
        text = clipped.rstrip(" .,;:") + "..."
    return text


def _webpage_preview(text: str, *, max_chars: int = WEBPAGE_MAX_PREVIEW_CHARS) -> str:
    out = " ".join(str(text or "").split()).strip()
    if len(out) <= max_chars:
        return out
    clipped = out[:max_chars]
    if " " in clipped[300:]:
        clipped = clipped[: clipped.rfind(" ")]
    return clipped.rstrip(" .,;:") + "..."


def _webpage_looks_bot_blocked(*, status_code: int, title: str, description: str, preview: str) -> bool:
    if int(status_code or 0) in WEBPAGE_BLOCK_STATUS_CODES:
        return True
    text = " ".join(
        [
            str(title or "").strip().lower(),
            str(description or "").strip().lower(),
            str(preview or "").strip().lower(),
        ]
    ).strip()
    if not text:
        return False
    return any(marker in text for marker in WEBPAGE_BOT_BLOCK_MARKERS)


def _fetch_webpage(
    normalized_url: str,
    *,
    timeout_sec: int,
    max_bytes: int,
) -> Dict[str, Any]:
    timeout_val = _coerce_int(timeout_sec, default=WEBPAGE_TIMEOUT_SEC, min_value=3, max_value=60)
    byte_limit = _coerce_int(max_bytes, default=WEBPAGE_MAX_RESPONSE_BYTES, min_value=65_536, max_value=25_000_000)

    attempts: List[Dict[str, Any]] = []
    selected: Optional[Dict[str, Any]] = None
    last_error = ""

    for user_agent in WEBPAGE_USER_AGENTS:
        session = requests.Session()
        response = None
        attempt: Dict[str, Any] = {
            "user_agent": user_agent,
            "status_code": 0,
            "blocked": False,
            "error": "",
        }
        try:
            response = session.get(
                normalized_url,
                headers=_webpage_request_headers(user_agent),
                timeout=(8, timeout_val),
                allow_redirects=True,
                stream=True,
            )
            attempt["status_code"] = int(response.status_code or 0)
            content_type = str(response.headers.get("Content-Type") or "")
            final_url = str(response.url or normalized_url)

            chunks: List[bytes] = []
            total = 0
            truncated = False
            for chunk in response.iter_content(chunk_size=64 * 1024):
                if not chunk:
                    continue
                next_total = total + len(chunk)
                if next_total > byte_limit:
                    keep = byte_limit - total
                    if keep > 0:
                        chunks.append(bytes(chunk[:keep]))
                        total += keep
                    truncated = True
                    break
                chunks.append(bytes(chunk))
                total = next_total
            raw = b"".join(chunks)

            if not _webpage_is_textual_content_type(content_type):
                selected = {
                    "ok": False,
                    "error": f"Non-text content type ({content_type or 'unknown'}). Use download_file instead.",
                    "url": final_url,
                    "content_type": content_type,
                    "status_code": int(response.status_code or 0),
                    "bytes": int(total),
                    "truncated": bool(truncated),
                    "attempts": attempts + [attempt],
                }
                break

            text = _webpage_decode_text(
                raw,
                content_type=content_type,
                response_encoding=str(response.encoding or ""),
            )
            text_preview = _webpage_preview(text)
            blocked = _webpage_looks_bot_blocked(
                status_code=int(response.status_code or 0),
                title="",
                description="",
                preview=text_preview,
            )
            attempt["blocked"] = bool(blocked)

            row = {
                "ok": True,
                "url": final_url,
                "content_type": content_type,
                "status_code": int(response.status_code or 0),
                "bytes": int(total),
                "truncated": bool(truncated),
                "raw_text": text,
                "attempts": attempts + [attempt],
            }

            if blocked:
                selected = row
                attempts.append(attempt)
                continue

            selected = row
            break
        except requests.RequestException as exc:
            message = str(exc).strip() or exc.__class__.__name__
            last_error = message
            attempt["error"] = message
            attempts.append(attempt)
            continue
        except Exception as exc:
            message = str(exc).strip() or exc.__class__.__name__
            last_error = message
            attempt["error"] = message
            attempts.append(attempt)
            continue
        finally:
            try:
                if response is not None:
                    response.close()
            except Exception:
                pass
            try:
                session.close()
            except Exception:
                pass

    if selected is None:
        return {
            "ok": False,
            "error": last_error or "Unable to fetch webpage.",
            "url": normalized_url,
            "attempts": attempts,
        }
    return selected


def _clean_redis_str(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (bytes, bytearray)):
        return value.decode("utf-8", errors="ignore").strip()
    return str(value).strip()


def _web_search_settings() -> Tuple[str, str]:
    api_key = ""
    cx = ""
    try:
        api_key = _clean_redis_str(redis_client.get(WEB_SEARCH_API_KEY_REDIS_KEY))
        cx = _clean_redis_str(redis_client.get(WEB_SEARCH_CX_REDIS_KEY))
    except Exception:
        api_key = ""
        cx = ""

    if api_key and cx:
        return api_key, cx

    try:
        legacy = redis_client.hgetall(WEB_SEARCH_LEGACY_SETTINGS_KEY) or {}
    except Exception:
        legacy = {}
    if not api_key:
        api_key = _clean_redis_str(legacy.get("GOOGLE_API_KEY") or legacy.get("google_api_key"))
    if not cx:
        cx = _clean_redis_str(legacy.get("GOOGLE_CX") or legacy.get("google_cx"))
    return api_key, cx


def search_web(
    query: str,
    *,
    num_results: int = 5,
    start: int = 1,
    site: Optional[str] = None,
    safe: str = "active",
    country: Optional[str] = None,
    language: Optional[str] = None,
    timeout_sec: int = WEB_SEARCH_TIMEOUT_SEC,
) -> Dict[str, Any]:
    q = str(query or "").strip()
    if not q:
        return {"tool": "search_web", "ok": False, "error": "query is required."}

    api_key, cx = _web_search_settings()
    if not api_key or not cx:
        return {
            "tool": "search_web",
            "ok": False,
            "error": "Web search is not configured. Set Google API Key and Search Engine ID (CX) in WebUI Settings.",
            "needs": [
                "Please set Google API Key in WebUI Settings > Web Search.",
                "Please set Google Search Engine ID (CX) in WebUI Settings > Web Search.",
            ],
        }

    max_results = _coerce_int(num_results, default=5, min_value=1, max_value=WEB_SEARCH_MAX_RESULTS)
    start_index = _coerce_int(start, default=1, min_value=1, max_value=91)
    timeout_val = _coerce_int(timeout_sec, default=WEB_SEARCH_TIMEOUT_SEC, min_value=3, max_value=60)
    safe_mode = str(safe or "active").strip().lower()
    if safe_mode not in {"active", "off"}:
        safe_mode = "active"

    params: Dict[str, Any] = {
        "key": api_key,
        "cx": cx,
        "q": q,
        "num": max_results,
        "start": start_index,
        "safe": safe_mode,
    }
    site_val = str(site or "").strip()
    if site_val:
        params["siteSearch"] = site_val
        params["siteSearchFilter"] = "i"

    country_val = str(country or "").strip().lower()
    if country_val and re.fullmatch(r"[a-z]{2}", country_val):
        params["gl"] = country_val

    language_val = str(language or "").strip().lower()
    if language_val:
        if language_val.startswith("lang_"):
            params["lr"] = language_val
        elif re.fullmatch(r"[a-z]{2}", language_val):
            params["lr"] = f"lang_{language_val}"

    endpoint = "https://www.googleapis.com/customsearch/v1"
    url = endpoint + "?" + urllib.parse.urlencode(params, doseq=True)
    req = urllib.request.Request(url, headers={"User-Agent": "Tater-AgentLab/1.0"})

    try:
        with urllib.request.urlopen(req, timeout=timeout_val) as resp:
            raw = resp.read(WEB_SEARCH_MAX_RESPONSE_BYTES)
    except urllib.error.HTTPError as e:
        body = ""
        try:
            body = e.read(1000).decode("utf-8", errors="replace")
        except Exception:
            body = ""
        message = body.strip() or str(e)
        return {"tool": "search_web", "ok": False, "error": f"Google CSE request failed ({e.code}): {message}"}
    except Exception as e:
        return {"tool": "search_web", "ok": False, "error": f"Web search failed: {e}"}

    try:
        payload = json.loads(raw.decode("utf-8", errors="replace"))
    except Exception:
        return {"tool": "search_web", "ok": False, "error": "Invalid response from Google CSE."}

    if isinstance(payload, dict) and payload.get("error"):
        err = payload.get("error") or {}
        msg = str(err.get("message") or "Unknown Google CSE error.")
        return {"tool": "search_web", "ok": False, "error": msg}

    items = payload.get("items") if isinstance(payload, dict) else []
    if not isinstance(items, list):
        items = []

    results: List[Dict[str, Any]] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        link = str(item.get("link") or "").strip()
        if not link:
            continue
        title = str(item.get("title") or "").strip() or link
        snippet = str(item.get("snippet") or "").strip()
        if len(snippet) > WEB_SEARCH_MAX_SNIPPET_CHARS:
            snippet = snippet[:WEB_SEARCH_MAX_SNIPPET_CHARS].rstrip() + "..."
        results.append(
            {
                "title": title,
                "url": link,
                "snippet": snippet,
                "display_url": str(item.get("displayLink") or "").strip(),
            }
        )

    search_info = payload.get("searchInformation") if isinstance(payload, dict) else {}
    search_time = None
    if isinstance(search_info, dict):
        search_time = search_info.get("searchTime")

    next_start = None
    total_results = None
    queries_blob = payload.get("queries") if isinstance(payload, dict) else None
    if isinstance(queries_blob, dict):
        req_pages = queries_blob.get("request")
        if isinstance(req_pages, list) and req_pages:
            req0 = req_pages[0] if isinstance(req_pages[0], dict) else {}
            raw_total = req0.get("totalResults")
            if raw_total is not None:
                try:
                    total_results = int(str(raw_total))
                except Exception:
                    total_results = None
        next_pages = queries_blob.get("nextPage")
        if isinstance(next_pages, list) and next_pages:
            np0 = next_pages[0] if isinstance(next_pages[0], dict) else {}
            raw_next = np0.get("startIndex")
            if raw_next is not None:
                try:
                    next_start = int(raw_next)
                except Exception:
                    next_start = None

    return {
        "tool": "search_web",
        "ok": True,
        "query": q,
        "start": start_index,
        "count": len(results),
        "num_results": max_results,
        "results": results,
        "site_filter": site_val or None,
        "search_time_sec": search_time,
        "total_results": total_results,
        "has_more": bool(next_start),
        "next_start": next_start,
    }


def _download_file_detect_media(path: Path, content_type: str) -> Tuple[str, str]:
    mime = str(content_type or "").split(";", 1)[0].strip().lower()
    if not mime:
        mime = str(mimetypes.guess_type(str(path))[0] or "").strip().lower()
    media_type = "file"
    if mime.startswith("image/"):
        media_type = "image"
    elif mime.startswith("audio/"):
        media_type = "audio"
    elif mime.startswith("video/"):
        media_type = "video"
    if not mime:
        if media_type == "image":
            mime = "image/png"
        elif media_type == "audio":
            mime = "audio/mpeg"
        elif media_type == "video":
            mime = "video/mp4"
        else:
            mime = "application/octet-stream"
    return media_type, mime


def _send_message_boolish(value: Any, default: bool = False) -> bool:
    if value is None:
        return bool(default)
    if isinstance(value, bool):
        return value
    text = str(value).strip().lower()
    if text in {"1", "true", "yes", "y", "on", "enabled"}:
        return True
    if text in {"0", "false", "no", "n", "off", "disabled"}:
        return False
    return bool(default)


def _send_message_load_settings() -> Dict[str, str]:
    return (
        redis_client.hgetall("verba_settings:Send Message")
        or redis_client.hgetall("verba_settings: Send Message")
        or {}
    )


def _send_message_last_macos_target() -> Dict[str, str]:
    out: Dict[str, str] = {}
    last_scope = str(redis_client.get("tater:macos:last_scope") or "").strip()
    last_device_id = str(redis_client.get("tater:macos:last_device_id") or "").strip()
    if last_scope:
        out["scope"] = last_scope
    if last_device_id:
        out["device_id"] = last_device_id
    return out


def _send_message_normalize_matrix_room_ref(room_ref: Any) -> str:
    ref = str(room_ref or "").strip()
    if not ref:
        return ""
    if ref.startswith("!") or ref.startswith("#"):
        return ref
    if ":" in ref:
        return f"#{ref}"
    return ref


def _send_message_extract_target_hint(raw: Any) -> str:
    text = str(raw or "").strip()
    if not text:
        return ""
    for pattern in (r"![^\s]+", r"#[A-Za-z0-9][A-Za-z0-9._:-]*", r"@[A-Za-z0-9_]+"):
        match = re.search(pattern, text)
        if match:
            return match.group(0)
    text = re.sub(r"^(?:room|channel|chat)\s+", "", text, flags=re.IGNORECASE)
    text = re.sub(
        r"\s+(?:to|in|on)\s+(?:discord|irc|matrix|telegram|home\s*assistant|homeassistant|ntfy|web\s*ui|webui)\b.*$",
        "",
        text,
        flags=re.IGNORECASE,
    )
    candidate = text.strip(" .")
    if not candidate:
        return ""
    if re.search(r"\s", candidate):
        return ""
    if normalize_notify_platform(candidate) in ALLOWED_PLATFORMS:
        return ""
    return candidate


def _send_message_coerce_targets(payload: Any) -> Dict[str, Any]:
    if isinstance(payload, dict):
        return dict(payload)
    if isinstance(payload, str):
        hint = _send_message_extract_target_hint(payload)
        if hint:
            return {"channel": hint}
    return {}


def _send_message_has_explicit_target(target_map: Dict[str, Any]) -> bool:
    for key in ("channel_id", "channel", "guild_id", "room_id", "room_alias", "chat_id", "scope", "device_id"):
        if str(target_map.get(key) or "").strip():
            return True
    return False


def _send_message_extract_destination_hint(raw: Any) -> str:
    text = str(raw or "").strip().lower()
    if not text:
        return ""
    match = re.search(
        r"\b(?:to|in|on)\s+(discord|irc|matrix|telegram|home\s*assistant|homeassistant|ntfy|web\s*ui|webui|mac\s*os|macos|my\s+mac)\b",
        text,
    )
    if not match:
        return ""
    token = str(match.group(1) or "").strip()
    return normalize_notify_platform(token)


def _send_message_extract_instruction_target(raw: Any) -> str:
    text = str(raw or "").strip()
    if not text:
        return ""
    match = re.search(r"\bto\s+([#!@][A-Za-z0-9][A-Za-z0-9._:-]*)", text, flags=re.IGNORECASE)
    if match:
        return str(match.group(1) or "").strip()
    return _send_message_extract_target_hint(text)


def _send_message_target_hint_from_map(target_map: Dict[str, Any]) -> str:
    for key in ("channel", "room_alias", "room_id", "chat_id", "channel_id"):
        token = str(target_map.get(key) or "").strip()
        if not token:
            continue
        extracted = _send_message_extract_target_hint(token)
        if extracted:
            return extracted
        return token
    return ""


def _send_message_target_variants(value: Any) -> set[str]:
    token = str(value or "").strip().lower().strip(" .,!?:;\"'")
    if not token:
        return set()
    token = re.sub(r"^(?:room|channel|chat)\s+", "", token, flags=re.IGNORECASE).strip()
    if not token:
        return set()
    variants = {token}
    plain = token.lstrip("#@!")
    if plain:
        variants.add(plain)
        variants.add(f"#{plain}")
        variants.add(f"@{plain}")
        variants.add(f"!{plain}")
    return {item for item in variants if item}


def _send_message_catalog_platform_rows(*, platform: Optional[str], limit: int) -> List[Dict[str, Any]]:
    try:
        payload = notifier_destination_catalog(
            redis_client=redis_client,
            platform=platform,
            limit=max(1, min(500, int(limit))),
        )
    except Exception:
        return []
    rows = payload.get("platforms") if isinstance(payload, dict) else []
    if not isinstance(rows, list):
        return []
    out: List[Dict[str, Any]] = []
    for item in rows:
        if isinstance(item, dict):
            out.append(item)
    return out


def _send_message_resolve_catalog_destination(
    *,
    destination: str,
    target_map: Dict[str, Any],
    origin: Optional[Dict[str, Any]],
) -> Tuple[str, Dict[str, Any], Optional[Dict[str, Any]]]:
    dest = normalize_notify_platform(destination)
    resolved_targets = dict(target_map or {})
    target_hint = _send_message_target_hint_from_map(resolved_targets)
    if not target_hint:
        return dest, resolved_targets, None
    hint_variants = _send_message_target_variants(target_hint)
    if not hint_variants:
        return dest, resolved_targets, None

    has_direct_id = False
    if dest == "discord":
        has_direct_id = bool(str(resolved_targets.get("channel_id") or "").strip())
    elif dest == "telegram":
        has_direct_id = bool(str(resolved_targets.get("chat_id") or "").strip())
    elif dest == "matrix":
        has_direct_id = bool(str(resolved_targets.get("room_id") or "").strip())

    catalog_rows = _send_message_catalog_platform_rows(platform=None, limit=200)
    if not catalog_rows:
        if dest in {"discord", "telegram", "matrix"} and target_hint and not has_direct_id:
            return dest, resolved_targets, {
                "unresolved": True,
                "reason": "destination_catalog_empty",
                "matched_hint": str(target_hint),
                "platform": str(dest),
                "choices": [],
            }
        return dest, resolved_targets, None

    # Prefer live/discovered and historical destinations over queued items.
    # Recent queue entries may include stale or failed attempts.
    source_rank = {"room_label": 8, "recent_history": 7, "default": 6, "recent_queue": 1}
    field_rank = {"channel_id": 7, "chat_id": 7, "room_id": 7, "room_alias": 6, "channel": 5}
    platform_priority = {
        "discord": 9,
        "matrix": 8,
        "irc": 7,
        "telegram": 6,
        "macos": 5,
        "homeassistant": 4,
        "ntfy": 3,
        "wordpress": 2,
        "webui": 1,
    }
    origin_platform = normalize_notify_platform(origin.get("platform")) if isinstance(origin, dict) else ""

    matches: List[Tuple[int, str, Dict[str, Any], str]] = []
    for platform_row in catalog_rows:
        platform_token = normalize_notify_platform(platform_row.get("platform"))
        if not platform_token:
            continue
        if dest and platform_token != dest:
            continue

        destinations = platform_row.get("destinations") if isinstance(platform_row, dict) else []
        if not isinstance(destinations, list):
            continue
        for row in destinations:
            if not isinstance(row, dict):
                continue
            targets = row.get("targets") if isinstance(row.get("targets"), dict) else {}
            matched_field = ""
            for field_name in ("channel_id", "channel", "room_id", "room_alias", "chat_id"):
                value = str(targets.get(field_name) or "").strip()
                if not value:
                    continue
                if hint_variants & _send_message_target_variants(value):
                    matched_field = field_name
                    break
            if not matched_field:
                continue

            score = int(field_rank.get(matched_field, 0))
            source = str(row.get("source") or "").strip().lower()
            score += int(source_rank.get(source, 0))
            score += int(platform_priority.get(platform_token, 0))
            if dest and platform_token == dest:
                score += 30
            if not dest and origin_platform and platform_token == origin_platform:
                score += 15
            matches.append((score, platform_token, row, matched_field))

    if not matches:
        if dest and target_hint and dest in ALLOWED_PLATFORMS:
            choices: List[Dict[str, Any]] = []
            platform_rows = _send_message_catalog_platform_rows(platform=dest, limit=80)
            if platform_rows:
                destinations = platform_rows[0].get("destinations") if isinstance(platform_rows[0], dict) else []
                if isinstance(destinations, list):
                    for row in destinations[:8]:
                        if not isinstance(row, dict):
                            continue
                        row_targets = row.get("targets") if isinstance(row.get("targets"), dict) else {}
                        choices.append(
                            {
                                "label": str(row.get("label") or "").strip(),
                                "targets": dict(row_targets),
                            }
                        )
            return dest, resolved_targets, {
                "unresolved": True,
                "reason": "target_not_found_in_catalog",
                "matched_hint": str(target_hint),
                "platform": str(dest),
                "choices": choices,
            }
        return dest, resolved_targets, None

    explicit_discord_channel_id = str(resolved_targets.get("channel_id") or "").strip()
    explicit_discord_guild_id = str(resolved_targets.get("guild_id") or "").strip()
    if not explicit_discord_channel_id and not explicit_discord_guild_id and (not dest or dest in {"discord", "webui"}):
        discord_rows_by_channel_id: Dict[str, Dict[str, Any]] = {}
        discord_rows_by_guild: Dict[str, Dict[str, Any]] = {}
        for score, platform_token, row, matched_field in matches:
            del score, matched_field
            if platform_token != "discord":
                continue
            targets = row.get("targets") if isinstance(row.get("targets"), dict) else {}
            channel_name = str(targets.get("channel") or "").strip()
            guild_token = str(targets.get("guild_id") or "").strip()
            channel_id_token = str(targets.get("channel_id") or "").strip()
            if not channel_name or not guild_token:
                if not channel_name or not channel_id_token:
                    continue
            if not (_send_message_target_variants(channel_name) & hint_variants):
                continue
            if channel_id_token and channel_id_token not in discord_rows_by_channel_id:
                discord_rows_by_channel_id[channel_id_token] = row
            if guild_token and guild_token not in discord_rows_by_guild:
                discord_rows_by_guild[guild_token] = row

        ambiguous = len(discord_rows_by_channel_id) > 1 or len(discord_rows_by_guild) > 1
        if ambiguous:
            source_rows: Dict[str, Dict[str, Any]]
            source_rows = discord_rows_by_channel_id if len(discord_rows_by_channel_id) > 1 else discord_rows_by_guild
            choices: List[Dict[str, Any]] = []
            for key_token, row in sorted(
                source_rows.items(),
                key=lambda item: str((item[1] or {}).get("label") or item[0]).lower(),
            ):
                targets = row.get("targets") if isinstance(row.get("targets"), dict) else {}
                choices.append(
                    {
                        "label": str(row.get("label") or key_token).strip() or key_token,
                        "targets": dict(targets),
                    }
                )
                if len(choices) >= 8:
                    break
            return "discord", resolved_targets, {
                "ambiguous": True,
                "reason": "discord_channel_name_ambiguous",
                "matched_hint": str(target_hint),
                "platform": "discord",
                "choices": choices,
            }

    matches.sort(
        key=lambda item: (
            -int(item[0]),
            str(item[1]),
            str((item[2] or {}).get("label") or ""),
        )
    )
    best_score, best_platform, best_row, best_field = matches[0]
    del best_score, best_field

    if dest and best_platform != dest:
        return dest, resolved_targets, None
    if not dest:
        dest = best_platform

    best_targets = best_row.get("targets") if isinstance(best_row.get("targets"), dict) else {}
    for key, value in best_targets.items():
        if str(resolved_targets.get(key) or "").strip():
            continue
        resolved_targets[key] = value

    if dest == "matrix":
        if not resolved_targets.get("room_id"):
            alias = resolved_targets.get("room_alias") or resolved_targets.get("channel")
            if alias:
                resolved_targets["room_id"] = _send_message_normalize_matrix_room_ref(alias)
        if resolved_targets.get("room_id"):
            resolved_targets["room_id"] = _send_message_normalize_matrix_room_ref(resolved_targets.get("room_id"))
    elif dest == "telegram":
        if not resolved_targets.get("chat_id"):
            chat_hint = resolved_targets.get("channel_id") or resolved_targets.get("channel")
            if chat_hint:
                resolved_targets["chat_id"] = chat_hint

    resolution = {
        "matched_hint": str(target_hint),
        "platform": str(dest),
        "label": str(best_row.get("label") or "").strip(),
        "source": str(best_row.get("source") or "").strip(),
        "targets": dict(best_targets),
    }
    return dest, resolved_targets, resolution


def _send_message_clean_attachment_payload(payload: Any) -> List[Dict[str, Any]]:
    if not isinstance(payload, list):
        return []
    out: List[Dict[str, Any]] = []
    for item in payload:
        if isinstance(item, dict):
            out.append(dict(item))
    return out


def _send_message_attachment_has_media_payload(item: Dict[str, Any]) -> bool:
    if not isinstance(item, dict):
        return False
    for key in ("bytes", "data", "path", "blob_key", "url"):
        value = item.get(key)
        if isinstance(value, (bytes, bytearray)) and value:
            return True
        if isinstance(value, str) and value.strip():
            return True
    return False


def _send_message_build_attachment_from_payload(
    payload: Any,
    *,
    fallback: Optional[Dict[str, Any]] = None,
) -> Optional[Dict[str, Any]]:
    if not isinstance(payload, dict):
        return None
    binary, filename, mimetype, _ = _read_artifact_bytes(payload)
    if binary is None:
        return None
    fallback_map = dict(fallback or {})
    fallback_name = str(fallback_map.get("name") or "").strip()
    final_name = str(fallback_name or filename or _artifact_name_from_path(payload.get("path")) or "file.bin").strip() or "file.bin"
    final_mime = _artifact_mimetype(final_name, fallback_map.get("mimetype") or mimetype or payload.get("mimetype"))
    final_type = _artifact_type(final_name, final_mime, fallback_map.get("type") or payload.get("type"))
    out: Dict[str, Any] = {
        "type": final_type,
        "name": final_name,
        "mimetype": final_mime,
        "bytes": binary,
        "size": len(binary),
    }
    artifact_id = str(fallback_map.get("artifact_id") or payload.get("artifact_id") or "").strip()
    if artifact_id:
        out["artifact_id"] = artifact_id
    return out


def _send_message_materialize_attachment(
    attachment: Any,
    *,
    origin: Optional[Dict[str, Any]],
) -> Optional[Dict[str, Any]]:
    if not isinstance(attachment, dict):
        return None
    item = dict(attachment)
    if _send_message_attachment_has_media_payload(item):
        name = str(item.get("name") or "file.bin").strip() or "file.bin"
        mimetype = _artifact_mimetype(name, item.get("mimetype"))
        kind = _artifact_type(name, mimetype, item.get("type"))
        out: Dict[str, Any] = {
            "type": kind,
            "name": name,
            "mimetype": mimetype,
        }
        for key in ("bytes", "data", "path", "blob_key", "url", "size"):
            value = item.get(key)
            if value in (None, ""):
                continue
            out[key] = value
        artifact_id = str(item.get("artifact_id") or "").strip()
        if artifact_id:
            out["artifact_id"] = artifact_id
        return out

    artifact_payload: Optional[Dict[str, Any]] = None
    artifact_id = str(item.get("artifact_id") or "").strip()
    if artifact_id:
        artifact_payload = _find_available_artifact(origin=origin, artifact_id=artifact_id)

    if artifact_payload is None:
        file_id = str(item.get("file_id") or item.get("id") or "").strip()
        if file_id:
            artifact_payload = {
                "file_id": file_id,
                "name": item.get("name"),
                "mimetype": item.get("mimetype"),
                "type": item.get("type"),
                "artifact_id": artifact_id,
            }

    if artifact_payload is None:
        name_hint = str(item.get("name") or "").strip().lower()
        if name_hint:
            for candidate in reversed(_origin_available_artifacts(origin)):
                if str(candidate.get("name") or "").strip().lower() == name_hint:
                    artifact_payload = candidate
                    break

    if artifact_payload is None:
        return None
    return _send_message_build_attachment_from_payload(
        artifact_payload,
        fallback=item,
    )


def _send_message_attachment_requested(raw: Any) -> bool:
    text = str(raw or "").strip().lower()
    if not text:
        return False
    if re.search(r"\b(?:this|that|the)\s+(?:image|photo|picture|screenshot|file|attachment)\b", text):
        return True
    if re.search(r"\b(?:attach|upload|share|send|post)\b", text) and re.search(
        r"\b(?:image|photo|picture|screenshot|file|attachment|document|doc)\b",
        text,
    ):
        return True
    return False


def _send_message_artifact_is_image(payload: Any) -> bool:
    if not isinstance(payload, dict):
        return False
    artifact_type = str(payload.get("type") or "").strip().lower()
    if artifact_type == "image":
        return True
    mimetype = str(payload.get("mimetype") or "").strip().lower()
    return mimetype.startswith("image/")


def _send_message_select_auto_attachment_payload(
    *,
    origin: Optional[Dict[str, Any]],
    request_text: str,
) -> Optional[Dict[str, Any]]:
    requested = _send_message_attachment_requested(request_text)
    current_turn = []
    if isinstance(origin, dict):
        raw_input = origin.get("input_artifacts")
        if isinstance(raw_input, list):
            current_turn = [dict(item) for item in raw_input if isinstance(item, dict)]

    available = [dict(item) for item in _origin_available_artifacts(origin) if isinstance(item, dict)]
    if not requested and len(current_turn) == 1:
        return current_turn[0]
    if not requested:
        return None

    candidate_pool = list(current_turn) + list(available)
    if not candidate_pool:
        return None

    text = str(request_text or "").strip().lower()
    hinted_id_match = re.search(r"\b(att\d+|a\d+)\b", text)
    if hinted_id_match:
        hinted_id = str(hinted_id_match.group(1) or "").strip().lower()
        for item in reversed(candidate_pool):
            if str(item.get("artifact_id") or "").strip().lower() == hinted_id:
                return item

    for item in reversed(candidate_pool):
        name = str(item.get("name") or "").strip().lower()
        if name and name in text:
            return item

    if re.search(r"\b(?:image|photo|picture|screenshot)\b", text):
        for item in reversed(candidate_pool):
            if _send_message_artifact_is_image(item):
                return item

    if len(current_turn) == 1:
        return current_turn[0]
    if len(candidate_pool) == 1:
        return candidate_pool[0]

    if "this " in text or text.startswith("this"):
        return candidate_pool[-1]

    return None


def _send_message_prepare_attachments(
    *,
    attachments: List[Dict[str, Any]],
    origin: Optional[Dict[str, Any]],
    request_text: str,
) -> Tuple[List[Dict[str, Any]], Dict[str, Any]]:
    resolved: List[Dict[str, Any]] = []
    resolved_refs = 0
    dropped_refs = 0

    for item in attachments:
        normalized = _send_message_materialize_attachment(item, origin=origin)
        if not isinstance(normalized, dict):
            dropped_refs += 1
            continue
        if str(normalized.get("artifact_id") or "").strip():
            resolved_refs += 1
        resolved.append(normalized)

    auto_added = False
    auto_name = ""
    auto_id = ""
    if not resolved:
        auto_payload = _send_message_select_auto_attachment_payload(
            origin=origin,
            request_text=request_text,
        )
        if isinstance(auto_payload, dict):
            auto_attachment = _send_message_build_attachment_from_payload(
                auto_payload,
                fallback=auto_payload,
            )
            if isinstance(auto_attachment, dict):
                auto_added = True
                auto_name = str(auto_attachment.get("name") or "").strip()
                auto_id = str(auto_attachment.get("artifact_id") or auto_payload.get("artifact_id") or "").strip()
                resolved.append(auto_attachment)

    meta = {
        "resolved_attachment_refs": int(resolved_refs),
        "dropped_attachment_refs": int(dropped_refs),
        "auto_attachment_added": bool(auto_added),
        "auto_attachment_name": auto_name,
        "auto_attachment_id": auto_id,
    }
    return resolved, meta


def _send_message_attachment_kind(mimetype: Any, fallback_type: Any = None) -> str:
    raw_type = str(fallback_type or "").strip().lower()
    if raw_type in {"image", "audio", "video", "file"}:
        return raw_type
    mime = str(mimetype or "").strip().lower()
    if mime.startswith("image/"):
        return "image"
    if mime.startswith("audio/"):
        return "audio"
    if mime.startswith("video/"):
        return "video"
    return "file"


def _artifact_name_from_path(path: Any) -> str:
    raw = str(path or "").strip().replace("\\", "/")
    if not raw:
        return ""
    return raw.rsplit("/", 1)[-1].strip()


def _artifact_mimetype(name: Any, mimetype: Any = "") -> str:
    mime = str(mimetype or "").strip().lower()
    if mime:
        return mime
    guessed = str(mimetypes.guess_type(str(name or "").strip())[0] or "").strip().lower()
    if guessed:
        return guessed
    return "application/octet-stream"


def _artifact_type(name: Any, mimetype: Any = "", fallback_type: Any = None) -> str:
    return _send_message_attachment_kind(_artifact_mimetype(name, mimetype), fallback_type)


def _origin_available_artifacts(origin: Optional[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not isinstance(origin, dict):
        return []
    raw = origin.get("available_artifacts")
    if not isinstance(raw, list):
        return []
    out: List[Dict[str, Any]] = []
    for item in raw:
        if isinstance(item, dict):
            out.append(dict(item))
    return out


def _find_available_artifact(
    *,
    origin: Optional[Dict[str, Any]],
    artifact_id: Any,
) -> Optional[Dict[str, Any]]:
    target = str(artifact_id or "").strip()
    if not target:
        return None
    for item in _origin_available_artifacts(origin):
        if str(item.get("artifact_id") or "").strip() == target:
            return item
    return None


def _read_artifact_bytes(
    payload: Any,
) -> Tuple[Optional[bytes], Optional[str], Optional[str], str]:
    if not isinstance(payload, dict):
        return None, None, None, "Artifact payload is invalid."

    if isinstance(payload.get("bytes"), (bytes, bytearray)):
        name = str(payload.get("name") or "file.bin").strip() or "file.bin"
        mime = _artifact_mimetype(name, payload.get("mimetype"))
        return bytes(payload.get("bytes")), name, mime, ""

    if isinstance(payload.get("data"), (bytes, bytearray)):
        name = str(payload.get("name") or "file.bin").strip() or "file.bin"
        mime = _artifact_mimetype(name, payload.get("mimetype"))
        return bytes(payload.get("data")), name, mime, ""

    if isinstance(payload.get("data"), str):
        decoded = _image_describe_decode_base64_payload(payload.get("data"))
        if decoded:
            name = str(payload.get("name") or "file.bin").strip() or "file.bin"
            mime = _artifact_mimetype(name, payload.get("mimetype"))
            return decoded, name, mime, ""

    path_value = str(payload.get("path") or "").strip()
    if path_value:
        resolved = _resolve_safe_path(path_value, [AGENT_LAB_DIR])
        if resolved is None:
            return None, None, None, "File path is outside the allowed workspace root."
        if not resolved.exists() or not resolved.is_file():
            return None, None, None, "File path does not exist."
        try:
            raw = resolved.read_bytes()
        except Exception:
            return None, None, None, "Failed to read the file."
        if not raw:
            return None, None, None, "The file is empty."
        name = str(payload.get("name") or resolved.name or "file.bin").strip() or "file.bin"
        mime = _artifact_mimetype(name, payload.get("mimetype"))
        return raw, name, mime, ""

    blob_client = _image_describe_blob_client()
    blob = _image_describe_load_blob_bytes(
        blob_client,
        blob_key=payload.get("blob_key"),
        file_id=payload.get("file_id") or payload.get("id"),
    )
    if blob:
        name = str(payload.get("name") or "file.bin").strip() or "file.bin"
        mime = _artifact_mimetype(name, payload.get("mimetype"))
        return blob, name, mime, ""

    return None, None, None, "Artifact bytes are unavailable."


def send_message(
    *,
    message: Any = None,
    content: Any = None,
    title: Any = None,
    platform: Any = None,
    targets: Any = None,
    attachments: Any = None,
    priority: Any = None,
    tags: Any = None,
    ttl_sec: Any = None,
    origin: Optional[Dict[str, Any]] = None,
    channel_id: Any = None,
    channel: Any = None,
    guild_id: Any = None,
    room_id: Any = None,
    room_alias: Any = None,
    device_service: Any = None,
    persistent: Any = None,
    api_notification: Any = None,
    chat_id: Any = None,
    device_id: Any = None,
    scope: Any = None,
) -> Dict[str, Any]:
    text_message = str(message or content or "").strip()
    destination = normalize_notify_platform(platform)
    target_map = _send_message_coerce_targets(targets)
    catalog_resolution: Optional[Dict[str, Any]] = None
    for key, value in (
        ("channel_id", channel_id),
        ("channel", channel),
        ("guild_id", guild_id),
        ("room_id", room_id),
        ("room_alias", room_alias),
        ("device_service", device_service),
        ("persistent", persistent),
        ("api_notification", api_notification),
        ("chat_id", chat_id),
        ("device_id", device_id),
        ("scope", scope),
    ):
        if value not in (None, "") and key not in target_map:
            target_map[key] = value

    attachment_items_raw = _send_message_clean_attachment_payload(attachments)
    attachment_items, attachment_meta = _send_message_prepare_attachments(
        attachments=attachment_items_raw,
        origin=origin,
        request_text=(message or content or text_message),
    )

    if not text_message and not attachment_items:
        return action_failure(
            code="missing_message",
            message="Cannot queue: missing message",
            needs=["Provide a message or include attachments to send."],
            say_hint="Ask for message content or an attachment to send.",
        )
    if not text_message and attachment_items:
        text_message = "Attachment"

    if not destination:
        destination_hint = _send_message_extract_destination_hint(message or content or text_message)
        if destination_hint:
            destination = destination_hint

    if not _send_message_has_explicit_target(target_map):
        inline_target_hint = _send_message_extract_instruction_target(message or content or text_message)
        if inline_target_hint:
            target_map["channel"] = inline_target_hint

    destination, target_map, catalog_resolution = _send_message_resolve_catalog_destination(
        destination=destination,
        target_map=target_map,
        origin=origin,
    )
    if isinstance(catalog_resolution, dict) and bool(catalog_resolution.get("ambiguous")):
        matched_hint = str(catalog_resolution.get("matched_hint") or target_map.get("channel") or "that channel").strip()
        choice_labels: List[str] = []
        for item in list(catalog_resolution.get("choices") or [])[:5]:
            if not isinstance(item, dict):
                continue
            label = str(item.get("label") or "").strip()
            if label:
                choice_labels.append(label)
        needs = ["Specify targets.guild_id or targets.channel_id for the destination Discord channel."]
        if choice_labels:
            needs.append("Matching destinations: " + " | ".join(choice_labels))
        return action_failure(
            code="ambiguous_discord_target",
            message=f"Cannot queue: `{matched_hint}` matches multiple Discord channels across different guilds.",
            needs=needs,
            say_hint="Ask the user which Discord guild/channel to use before sending.",
        )
    if isinstance(catalog_resolution, dict) and bool(catalog_resolution.get("unresolved")):
        unresolved_platform = normalize_notify_platform(catalog_resolution.get("platform") or destination)
        if unresolved_platform in {"discord", "matrix", "telegram"}:
            matched_hint = str(catalog_resolution.get("matched_hint") or "").strip() or "that destination"
            choice_labels: List[str] = []
            for item in list(catalog_resolution.get("choices") or [])[:5]:
                if not isinstance(item, dict):
                    continue
                label = str(item.get("label") or "").strip()
                if label:
                    choice_labels.append(label)
            needs = [
                f"Provide an exact {unresolved_platform} target id/name for `{matched_hint}` (for example channel_id/chat_id/room_id).",
            ]
            if choice_labels:
                needs.append("Known destinations: " + " | ".join(choice_labels))
            return action_failure(
                code="unknown_destination_target",
                message=f"Cannot queue: `{matched_hint}` was not found in known {unresolved_platform} destinations.",
                needs=needs,
                say_hint="Ask the user to pick a known destination or provide an exact destination id.",
            )

    if not destination and isinstance(origin, dict):
        origin_platform = normalize_notify_platform(origin.get("platform"))
        if origin_platform in ALLOWED_PLATFORMS:
            destination = origin_platform

    if destination == "macos" and not target_map.get("scope") and not target_map.get("device_id"):
        inferred_target = _send_message_last_macos_target()
        if inferred_target.get("scope"):
            target_map["scope"] = inferred_target["scope"]
        if inferred_target.get("device_id"):
            target_map["device_id"] = inferred_target["device_id"]

    if destination not in ALLOWED_PLATFORMS:
        return action_failure(
            code="missing_destination_platform",
            message="Cannot queue: missing destination platform",
            needs=[
                "Specify a destination platform such as discord, matrix, telegram, macos, homeassistant, ntfy, irc, or webui.",
                "For macOS you can also say 'my mac' or 'mac os'.",
            ],
            say_hint="Explain that a destination platform is required.",
        )

    origin_platform = normalize_notify_platform(origin.get("platform")) if isinstance(origin, dict) else ""
    origin_has_macos_target = (
        origin_platform == "macos"
        and bool(str(origin.get("scope") or "").strip() or str(origin.get("device_id") or "").strip())
    ) if isinstance(origin, dict) else False
    if destination == "macos" and not target_map.get("scope") and not target_map.get("device_id") and not origin_has_macos_target:
        return action_failure(
            code="missing_macos_target",
            message="Cannot queue: missing target device/scope",
            needs=[
                "Say 'to my mac' to target the last active Mac device, or specify targets.scope / targets.device_id explicitly.",
                "If this is your first macOS message from WebUI, open the Mac app once so Tater can learn your current device.",
            ],
            say_hint="Ask which Mac device to target, and mention 'my mac' as shorthand.",
        )

    if destination == "matrix":
        if not target_map.get("room_id"):
            alias = target_map.get("room_alias") or target_map.get("channel")
            if alias:
                target_map["room_id"] = alias
        if target_map.get("room_id"):
            target_map["room_id"] = _send_message_normalize_matrix_room_ref(target_map.get("room_id"))
    elif destination == "homeassistant":
        if "api_notification" not in target_map:
            settings = _send_message_load_settings()
            target_map["api_notification"] = _send_message_boolish(settings.get("ENABLE_HA_API_NOTIFICATION"), True)

    meta = {
        "priority": priority,
        "tags": tags,
        "ttl_sec": ttl_sec,
    }

    result = dispatch_notification_sync(
        platform=destination,
        title=str(title or "").strip() or None,
        content=text_message,
        targets=target_map,
        origin=origin,
        meta=meta,
        attachments=attachment_items,
    )
    if str(result or "").strip().lower().startswith("cannot queue"):
        return action_failure(
            code="send_message_failed",
            message=str(result or "").strip() or "Failed to queue notification.",
            say_hint="Explain why the notification could not be queued.",
        )
    return action_success(
        facts={
            "platform": destination,
            "target_count": len([value for value in target_map.values() if value not in (None, "", False)]),
            "attachment_count": len(attachment_items),
            "catalog_target_resolved": bool(catalog_resolution),
            "resolved_attachment_refs": int(attachment_meta.get("resolved_attachment_refs") or 0),
            "auto_attachment_added": bool(attachment_meta.get("auto_attachment_added")),
        },
        data={
            "result": str(result or "").strip(),
            "platform": destination,
            "targets": target_map,
            "attachment_count": len(attachment_items),
            "catalog_resolution": (catalog_resolution or {}),
            "resolved_attachment_refs": int(attachment_meta.get("resolved_attachment_refs") or 0),
            "dropped_attachment_refs": int(attachment_meta.get("dropped_attachment_refs") or 0),
            "auto_attachment_added": bool(attachment_meta.get("auto_attachment_added")),
            "auto_attachment_name": str(attachment_meta.get("auto_attachment_name") or ""),
            "auto_attachment_id": str(attachment_meta.get("auto_attachment_id") or ""),
        },
        summary_for_user=str(result or "").strip(),
        say_hint="Confirm the queued notification destination and keep it brief.",
    )


def _attach_file_notify_suggestions(origin: Optional[Dict[str, Any]], *, limit: int = 5) -> List[Dict[str, Any]]:
    if not isinstance(origin, dict):
        return []
    platform = normalize_notify_platform(origin.get("platform"))
    if not platform or platform not in ALLOWED_PLATFORMS:
        return []

    catalog_rows = _send_message_catalog_platform_rows(platform=platform, limit=max(1, min(20, int(limit))))
    if not catalog_rows:
        return []
    platform_row = catalog_rows[0]
    destinations = platform_row.get("destinations") if isinstance(platform_row, dict) else []
    if not isinstance(destinations, list):
        return []

    out: List[Dict[str, Any]] = []
    for row in destinations[: max(1, min(20, int(limit)))]:
        if not isinstance(row, dict):
            continue
        targets = row.get("targets") if isinstance(row.get("targets"), dict) else {}
        out.append(
            {
                "platform": platform,
                "label": str(row.get("label") or platform).strip() or platform,
                "targets": dict(targets),
                "source": str(row.get("source") or "").strip(),
            }
        )
    return out


def attach_file(
    *,
    artifact_id: Any = None,
    path: Any = None,
    message: Any = None,
    content: Any = None,
    title: Any = None,
    platform: Any = None,
    targets: Any = None,
    priority: Any = None,
    tags: Any = None,
    ttl_sec: Any = None,
    channel_id: Any = None,
    channel: Any = None,
    guild_id: Any = None,
    room_id: Any = None,
    room_alias: Any = None,
    device_service: Any = None,
    persistent: Any = None,
    api_notification: Any = None,
    chat_id: Any = None,
    device_id: Any = None,
    scope: Any = None,
    origin: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    artifact_payload: Optional[Dict[str, Any]] = None
    chosen_artifact_id = str(artifact_id or "").strip()
    explicit_path = str(path or "").strip()

    if chosen_artifact_id:
        artifact_payload = _find_available_artifact(origin=origin, artifact_id=chosen_artifact_id)
        if artifact_payload is None:
            return action_failure(
                code="artifact_not_found",
                message=f"Artifact `{chosen_artifact_id}` was not found for this conversation.",
                needs=["Use an artifact_id from the available conversation artifacts or provide a file path."],
                say_hint="Explain that the requested artifact is unavailable in this conversation.",
            )
    elif explicit_path:
        artifact_payload = {"path": explicit_path}
    else:
        return action_failure(
            code="missing_artifact_reference",
            message="attach_file requires an artifact_id or path.",
            needs=["Provide an artifact_id from this conversation or a local file path to attach."],
            say_hint="Ask for the exact file to attach.",
        )

    binary, filename, mimetype, error_message = _read_artifact_bytes(artifact_payload)
    if binary is None:
        return action_failure(
            code="artifact_unavailable",
            message=error_message or "The requested file could not be attached.",
            needs=["Provide a valid artifact_id from this conversation or a readable local file path."],
            say_hint="Explain briefly why the file could not be attached.",
        )

    final_name = str(filename or _artifact_name_from_path(explicit_path) or "file.bin").strip() or "file.bin"
    final_mime = _artifact_mimetype(final_name, mimetype)
    final_type = _artifact_type(final_name, final_mime, artifact_payload.get("type") if isinstance(artifact_payload, dict) else None)
    artifact_out = {
        "type": final_type,
        "name": final_name,
        "mimetype": final_mime,
        "bytes": binary,
        "size": len(binary),
    }

    if chosen_artifact_id:
        artifact_out["artifact_id"] = chosen_artifact_id

    text_message = str(message or content or "").strip()
    destination = normalize_notify_platform(platform)
    target_map = _send_message_coerce_targets(targets)
    for key, value in (
        ("channel_id", channel_id),
        ("channel", channel),
        ("guild_id", guild_id),
        ("room_id", room_id),
        ("room_alias", room_alias),
        ("device_service", device_service),
        ("persistent", persistent),
        ("api_notification", api_notification),
        ("chat_id", chat_id),
        ("device_id", device_id),
        ("scope", scope),
    ):
        if value not in (None, "") and key not in target_map:
            target_map[key] = value

    if not destination:
        destination_hint = _send_message_extract_destination_hint(text_message)
        if destination_hint:
            destination = destination_hint
    if not _send_message_has_explicit_target(target_map):
        inline_target_hint = _send_message_extract_instruction_target(text_message)
        if inline_target_hint:
            target_map["channel"] = inline_target_hint

    has_delivery_intent = bool(destination or _send_message_has_explicit_target(target_map))
    notify_suggestions = _attach_file_notify_suggestions(origin, limit=5)
    if has_delivery_intent:
        try:
            send_result = send_message(
                message=(text_message or "Attachment"),
                title=title,
                platform=destination,
                targets=target_map,
                attachments=[artifact_out],
                priority=priority,
                tags=tags,
                ttl_sec=ttl_sec,
                origin=origin,
                channel_id=channel_id,
                channel=channel,
                guild_id=guild_id,
                room_id=room_id,
                room_alias=room_alias,
                device_service=device_service,
                persistent=persistent,
                api_notification=api_notification,
                chat_id=chat_id,
                device_id=device_id,
                scope=scope,
            )
        except Exception as exc:
            return action_failure(
                code="attach_file_send_failed",
                message=f"Failed to send attachment to destination: {str(exc) or 'unknown error'}",
                say_hint="Explain why the attachment could not be sent.",
            )
        if not isinstance(send_result, dict) or not bool(send_result.get("ok")):
            return send_result if isinstance(send_result, dict) else action_failure(
                code="attach_file_send_failed",
                message="Failed to send attachment to destination.",
                say_hint="Explain why the file could not be sent.",
            )

        send_data = send_result.get("data") if isinstance(send_result.get("data"), dict) else {}
        send_summary = str(send_result.get("summary_for_user") or "").strip() or "Queued attachment delivery."
        return action_success(
            facts={
                "artifact_id": chosen_artifact_id,
                "name": final_name,
                "size": len(binary),
                "delivery_requested": True,
                "platform": str(send_data.get("platform") or destination or "").strip(),
                "notify_suggestion_count": len(notify_suggestions),
            },
            data={
                "artifact_id": chosen_artifact_id,
                "name": final_name,
                "mimetype": final_mime,
                "size": len(binary),
                "notify_suggestions": notify_suggestions,
                "delivery": send_data,
            },
            summary_for_user=f"{send_summary} (attached {final_name})",
            say_hint="Confirm the destination and attached file in one short sentence.",
            artifacts=[artifact_out],
        )

    return action_success(
        facts={
            "artifact_id": chosen_artifact_id,
            "name": final_name,
            "size": len(binary),
            "delivery_requested": False,
            "notify_suggestion_count": len(notify_suggestions),
        },
        data={
            "artifact_id": chosen_artifact_id,
            "name": final_name,
            "mimetype": final_mime,
            "size": len(binary),
            "notify_suggestions": notify_suggestions,
        },
        summary_for_user=f"Attached {final_name}.",
        say_hint="Confirm the file attachment briefly.",
        artifacts=[artifact_out],
    )


def _image_describe_blob_client() -> redis.Redis:
    return redis_blob_client


def _image_describe_blob_key_candidates(*, blob_key: Any = None, file_id: Any = None) -> List[str]:
    out: List[str] = []
    blob = _as_text(blob_key).strip()
    if blob:
        out.append(blob)

    fid = _as_text(file_id).strip()
    if fid:
        if fid.startswith((WEBUI_FILE_BLOB_KEY_PREFIX, "tater:blob:", "tater:matrix:")):
            out.append(fid)
        else:
            out.append(f"{WEBUI_FILE_BLOB_KEY_PREFIX}{fid}")
            out.append(fid)

    unique: List[str] = []
    seen = set()
    for key in out:
        if key and key not in seen:
            unique.append(key)
            seen.add(key)
    return unique


def _image_describe_load_blob_bytes(
    blob_client: redis.Redis,
    *,
    blob_key: Any = None,
    file_id: Any = None,
) -> Optional[bytes]:
    for key in _image_describe_blob_key_candidates(blob_key=blob_key, file_id=file_id):
        try:
            data = blob_client.get(key)
        except Exception:
            data = None
        if data is None:
            continue
        if isinstance(data, (bytes, bytearray)):
            return bytes(data)
        if isinstance(data, str):
            return data.encode("utf-8", errors="replace")
    return None


def _image_describe_decode_base64_payload(data: Any) -> Optional[bytes]:
    text = _as_text(data).strip()
    if not text:
        return None
    if text.startswith("data:") and "," in text:
        text = text.split(",", 1)[1]
    pad = len(text) % 4
    if pad:
        text += "=" * (4 - pad)
    try:
        decoded = base64.b64decode(text)
    except Exception:
        return None
    return bytes(decoded) if decoded else None


def _image_describe_mime_allowed(mimetype: Any) -> bool:
    mime = _as_text(mimetype).strip().lower()
    if not mime:
        return False
    return mime in VISION_ALLOWED_MIMETYPES


def _image_describe_looks_like_http_url(value: Any) -> bool:
    text = _as_text(value).strip()
    if not text:
        return False
    parsed = urllib.parse.urlparse(text)
    return parsed.scheme.lower() in {"http", "https"} and bool(parsed.netloc)


def _image_describe_normalize_filename(name: Any, mimetype: Any = "") -> str:
    text = _as_text(name).strip()
    if text:
        return text
    mime = _as_text(mimetype).strip().lower()
    if mime in {"image/jpeg", "image/jpg"}:
        return "image.jpg"
    if mime == "image/webp":
        return "image.webp"
    if mime == "image/gif":
        return "image.gif"
    if mime == "image/bmp":
        return "image.bmp"
    if mime == "image/tiff":
        return "image.tiff"
    return "image.png"


def _image_describe_to_data_url(image_bytes: bytes, filename: str) -> str:
    mime = mimetypes.guess_type(filename or "")[0] or "image/png"
    b64 = base64.b64encode(image_bytes).decode("utf-8")
    return f"data:{mime};base64,{b64}"


def _image_describe_download_image_url(
    value: Any,
) -> Tuple[Optional[bytes], Optional[str], Optional[str], Optional[str]]:
    raw_url = _as_text(value).strip()
    if not raw_url:
        return None, None, None, "url_empty"
    if not _image_describe_looks_like_http_url(raw_url):
        return None, None, None, "url_invalid"

    try:
        with requests.get(
            raw_url,
            timeout=30,
            stream=True,
            allow_redirects=True,
            headers={"User-Agent": "Tater/1.0"},
        ) as response:
            if response.status_code >= 300:
                return None, None, None, "url_http_error"

            content_type = _as_text(response.headers.get("Content-Type") or "").split(";", 1)[0].strip().lower()
            chunks: List[bytes] = []
            for chunk in response.iter_content(chunk_size=65536):
                if not chunk:
                    continue
                chunks.append(chunk)
            data = b"".join(chunks)
            final_url = _as_text(response.url).strip() or raw_url
    except requests.RequestException:
        return None, None, None, "url_request_failed"
    except Exception:
        return None, None, None, "url_request_failed"

    if not data:
        return None, None, None, "url_empty_response"

    parsed = urllib.parse.urlparse(final_url)
    guessed_name = Path(parsed.path).name if parsed.path else ""
    filename = _image_describe_normalize_filename(guessed_name, content_type)
    guessed_mime = _as_text(mimetypes.guess_type(filename)[0]).strip().lower()
    mimetype = content_type or guessed_mime or "image/png"

    if content_type and not content_type.startswith("image/"):
        return None, None, None, "url_not_image"
    if not content_type and (not guessed_mime or not guessed_mime.startswith("image/")):
        return None, None, None, "url_not_image"
    if mimetype.startswith("image/") and not _image_describe_mime_allowed(mimetype):
        return None, None, None, "url_unsupported_type"

    return data, filename, mimetype, None


def _image_describe_extract_from_payload(
    blob_client: redis.Redis,
    payload: Any,
) -> Tuple[Optional[bytes], Optional[str], Optional[str]]:
    if payload is None:
        return None, None, None

    if isinstance(payload, dict) and payload.get("marker") == "plugin_response":
        return _image_describe_extract_from_payload(blob_client, payload.get("content"))

    if isinstance(payload, list):
        for item in payload:
            raw, name, mime = _image_describe_extract_from_payload(blob_client, item)
            if raw:
                return raw, name, mime
        return None, None, None

    if not isinstance(payload, dict):
        return None, None, None

    media_type = _as_text(payload.get("type")).strip().lower()
    mimetype = _as_text(payload.get("mimetype")).strip().lower()
    if not mimetype:
        guessed_name = _as_text(payload.get("name")).strip()
        mimetype = _as_text(mimetypes.guess_type(guessed_name)[0]).strip().lower()

    if media_type in {"image", "file"}:
        if media_type == "file" and (not mimetype or not mimetype.startswith("image/")):
            return None, None, None
        if mimetype and not _image_describe_mime_allowed(mimetype):
            return None, None, None

        filename = _image_describe_normalize_filename(payload.get("name"), mimetype)

        if isinstance(payload.get("bytes"), (bytes, bytearray)):
            return bytes(payload["bytes"]), filename, mimetype or "image/png"
        if isinstance(payload.get("data"), (bytes, bytearray)):
            return bytes(payload["data"]), filename, mimetype or "image/png"

        decoded = _image_describe_decode_base64_payload(payload.get("data"))
        if decoded:
            return decoded, filename, mimetype or "image/png"

        blob = _image_describe_load_blob_bytes(
            blob_client,
            blob_key=payload.get("blob_key"),
            file_id=payload.get("id") or payload.get("file_id"),
        )
        if blob:
            mm = mimetype or _as_text(mimetypes.guess_type(filename)[0]).strip().lower() or "image/png"
            return blob, filename, mm

        ref_url = payload.get("url")
        if _image_describe_looks_like_http_url(ref_url):
            raw, remote_name, remote_mime, err = _image_describe_download_image_url(ref_url)
            if raw and not err:
                final_name = _image_describe_normalize_filename(payload.get("name") or remote_name, mimetype or remote_mime)
                final_mime = (
                    mimetype
                    or remote_mime
                    or _as_text(mimetypes.guess_type(final_name)[0]).strip().lower()
                    or "image/png"
                )
                return raw, final_name, final_mime

        return None, None, None

    if payload.get("blob_key") or payload.get("file_id") or payload.get("id"):
        blob = _image_describe_load_blob_bytes(
            blob_client,
            blob_key=payload.get("blob_key"),
            file_id=payload.get("file_id") or payload.get("id"),
        )
        if blob:
            filename = _image_describe_normalize_filename(payload.get("name"), payload.get("mimetype"))
            mimetype = (
                _as_text(payload.get("mimetype")).strip().lower()
                or _as_text(mimetypes.guess_type(filename)[0]).strip().lower()
            )
            if not mimetype:
                mimetype = "image/png"
            if mimetype.startswith("image/"):
                return blob, filename, mimetype

    if _image_describe_looks_like_http_url(payload.get("url")):
        raw, filename, mimetype, err = _image_describe_download_image_url(payload.get("url"))
        if raw and not err:
            return raw, filename, mimetype

    return None, None, None


def _image_describe_resolve_explicit_image(
    *,
    prompt: Any = None,
    query: Any = None,
    request: Any = None,
    artifact_id: Any = None,
    url: Any = None,
    path: Any = None,
    blob_key: Any = None,
    file_id: Any = None,
    image_ref: Any = None,
    source: Any = None,
    file: Any = None,
    name: Any = None,
    mimetype: Any = None,
    origin: Optional[Dict[str, Any]] = None,
) -> Tuple[Optional[bytes], Optional[str], Optional[str], Optional[str], str]:
    del prompt, query, request
    blob_client = _image_describe_blob_client()

    artifact_token = _as_text(artifact_id).strip()
    if artifact_token:
        artifact_payload = _find_available_artifact(origin=origin, artifact_id=artifact_token)
        if artifact_payload is None:
            return None, None, None, "artifact_not_found", f"Artifact `{artifact_token}` was not found for this conversation."
        raw, filename, mime, err = _read_artifact_bytes(artifact_payload)
        if raw is None:
            return None, None, None, "artifact_unavailable", err or "The requested artifact could not be read."
        final_name = _image_describe_normalize_filename(filename or artifact_payload.get("name"), mime)
        final_mime = _as_text(mime).strip().lower() or _artifact_mimetype(final_name, artifact_payload.get("mimetype"))
        if not final_mime.startswith("image/"):
            return None, None, None, "artifact_not_image", "The selected artifact is not an image."
        if not _image_describe_mime_allowed(final_mime):
            return None, None, None, "artifact_unsupported_type", "The selected image type is not supported."
        return raw, final_name, final_mime, "artifact", ""

    for ref in (image_ref,):
        if isinstance(ref, dict):
            image_bytes, filename, mime = _image_describe_extract_from_payload(blob_client, ref)
            if image_bytes:
                return image_bytes, filename, mime, "explicit_ref", ""

    explicit_url = _as_text(url).strip()
    if explicit_url:
        if not _image_describe_looks_like_http_url(explicit_url):
            return None, None, None, "url_invalid", "Image URL must be a valid http/https URL."
    else:
        source_hint = source or file
        if _image_describe_looks_like_http_url(source_hint):
            explicit_url = _as_text(source_hint).strip()

    if explicit_url:
        data, filename, mime, err = _image_describe_download_image_url(explicit_url)
        if err:
            msg_map = {
                "url_empty": "Image URL is empty.",
                "url_invalid": "Image URL must be a valid http/https URL.",
                "url_request_failed": "Failed to download the image URL.",
                "url_http_error": "Image URL request returned an HTTP error.",
                "url_empty_response": "Image URL returned no data.",
                "url_not_image": "URL did not resolve to an image.",
                "url_unsupported_type": "Image URL returned an unsupported image type.",
            }
            return None, None, None, err, msg_map.get(err, "Invalid image URL.")
        if data:
            return data, filename, mime, "url", ""

    image_path = path
    if not _as_text(image_path).strip():
        source_hint = source or file
        if not _image_describe_looks_like_http_url(source_hint):
            image_path = source_hint
    if _as_text(image_path).strip():
        resolved = _resolve_safe_path(_as_text(image_path), [AGENT_LAB_DIR])
        if resolved is None:
            return None, None, None, "path_outside_workspace", "Image path is outside the allowed workspace root."
        if not resolved.exists() or not resolved.is_file():
            return None, None, None, "path_missing", "Image path does not exist."
        try:
            data = resolved.read_bytes()
        except Exception:
            return None, None, None, "path_read_failed", "Failed to read image from the provided path."
        if not data:
            return None, None, None, "path_empty", "The provided image file is empty."
        filename = resolved.name or "image.png"
        mime = _as_text(mimetypes.guess_type(filename)[0]).strip().lower() or "image/png"
        if mime and not mime.startswith("image/"):
            return None, None, None, "path_not_image", "The provided path is not an image file."
        return data, filename, mime, "path", ""

    blob = _image_describe_load_blob_bytes(
        blob_client,
        blob_key=blob_key,
        file_id=file_id,
    )
    if blob:
        filename = _image_describe_normalize_filename(name, mimetype)
        mime = (
            _as_text(mimetype).strip().lower()
            or _as_text(mimetypes.guess_type(filename)[0]).strip().lower()
            or "image/png"
        )
        if mime and not mime.startswith("image/"):
            return None, None, None, "blob_not_image", "The provided blob/file reference is not an image."
        return blob, filename, mime, "blob", ""

    return None, None, None, "", ""


def _image_describe_call_vision_api(
    *,
    image_bytes: bytes,
    filename: str,
    prompt: str,
) -> Tuple[Optional[str], Optional[str], Optional[str]]:
    settings = get_vision_settings(
        default_api_base=DEFAULT_VISION_API_BASE,
        default_model=DEFAULT_VISION_MODEL,
    )
    api_base = _as_text(settings.get("api_base")).strip().rstrip("/")
    model = _as_text(settings.get("model")).strip()
    api_key = _as_text(settings.get("api_key")).strip()

    if not api_base or not model:
        return None, None, "Vision settings are incomplete. Configure API base and model in Settings."

    url = f"{api_base}/v1/chat/completions"
    payload = {
        "model": model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": _as_text(prompt).strip() or VISION_DEFAULT_PROMPT},
                    {"type": "image_url", "image_url": {"url": _image_describe_to_data_url(image_bytes, filename)}},
                ],
            }
        ],
        "temperature": 0.2,
    }

    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    try:
        response = requests.post(url, json=payload, headers=headers, timeout=90)
    except Exception as exc:
        return None, model, f"Vision API request failed: {exc}"

    if response.status_code >= 300:
        detail = _as_text(response.text).strip()
        if detail:
            detail = detail[:400]
            return None, model, f"Vision API request failed with HTTP {response.status_code}: {detail}"
        return None, model, f"Vision API request failed with HTTP {response.status_code}."

    try:
        parsed = response.json()
    except Exception:
        return None, model, "Vision API returned non-JSON output."

    try:
        content = parsed["choices"][0]["message"]["content"]
    except Exception:
        return None, model, "Vision API response did not include a valid assistant message."

    description = ""
    if isinstance(content, str):
        description = content.strip()
    elif isinstance(content, list):
        chunks: List[str] = []
        for item in content:
            if isinstance(item, str) and item.strip():
                chunks.append(item.strip())
            elif isinstance(item, dict) and _as_text(item.get("text")).strip():
                chunks.append(_as_text(item.get("text")).strip())
        description = "\n".join(chunks).strip()

    if not description:
        return None, model, "Vision API returned an empty description."

    return description, model, None


def image_describe(
    *,
    request: Any = None,
    query: Any = None,
    prompt: Any = None,
    artifact_id: Any = None,
    url: Any = None,
    path: Any = None,
    blob_key: Any = None,
    file_id: Any = None,
    image_ref: Any = None,
    source: Any = None,
    file: Any = None,
    name: Any = None,
    mimetype: Any = None,
    platform: Optional[str] = None,
    origin: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    origin_payload = origin if isinstance(origin, dict) else {}
    origin_request = _as_text(
        origin_payload.get("request_text")
        or origin_payload.get("raw_message")
        or origin_payload.get("raw")
    ).strip()
    prompt_text = _as_text(prompt or query or request or origin_request).strip() or VISION_DEFAULT_PROMPT

    image_bytes, filename, mime, resolution_source, error_message = _image_describe_resolve_explicit_image(
        prompt=prompt,
        query=query,
        request=request,
        artifact_id=artifact_id,
        url=url,
        path=path,
        blob_key=blob_key,
        file_id=file_id,
        image_ref=image_ref,
        source=source,
        file=file,
        name=name,
        mimetype=mimetype,
        origin=origin,
    )

    if image_bytes is None and error_message:
        return action_failure(
            code="invalid_image_source",
            message=error_message,
            needs=[
                "Use an artifact_id from this conversation, provide an image URL, or provide /downloads/... or /documents/... path."
            ],
            say_hint="Ask for a valid image source and keep guidance brief.",
        )

    if image_bytes is None:
        return action_failure(
            code="no_image_found",
            message="No image was found. Use an artifact_id from this conversation, provide an image URL, or provide a path in /downloads or /documents.",
            needs=["Please provide an image source to describe."],
            say_hint="Ask for an image artifact_id, URL, or a path in /downloads or /documents.",
        )

    filename = _image_describe_normalize_filename(filename, mime)
    description, model, error = _image_describe_call_vision_api(
        image_bytes=image_bytes,
        filename=filename,
        prompt=prompt_text,
    )
    if error:
        return action_failure(
            code="vision_request_failed",
            message=error,
            say_hint="Explain the vision request failure and ask whether to retry.",
        )

    text = _as_text(description).strip()
    return action_success(
        facts={
            "tool": "image_describe",
            "source": resolution_source or "unknown",
            "filename": filename,
        },
        data={
            "description": text,
            "text": text,
            "filename": filename,
            "mimetype": mime or _as_text(mimetypes.guess_type(filename)[0]).strip() or "image/png",
            "model": model or "",
            "source": resolution_source or "unknown",
        },
        summary_for_user=text,
        say_hint="Return the image description directly and do not invent extra visual details.",
    )


def _ai_tasks_normalize_channel_targets(dest: str, targets: Dict[str, Any]) -> Dict[str, Any]:
    t = dict(targets or {})
    channel_ref = t.get("channel")

    if not channel_ref:
        if dest == "discord":
            channel_ref = t.get("channel_id")
        elif dest == "matrix":
            channel_ref = t.get("room_id")
        elif dest == "homeassistant":
            channel_ref = t.get("device_service")
        elif dest == "telegram":
            channel_ref = t.get("chat_id")

    if not channel_ref:
        return t

    ref = str(channel_ref).strip()
    if not ref:
        return {}

    if dest == "discord":
        if ref.isdigit():
            return {"channel_id": ref}
        return {"channel": ref}

    if dest == "matrix":
        if not ref.startswith(("!", "#")) and ":" in ref:
            ref = f"#{ref}"
        return {"room_id": ref}

    if dest == "homeassistant":
        return {"device_service": ref}

    if dest == "telegram":
        return {"chat_id": ref}

    return {"channel": ref}


def _ai_tasks_extract_time_of_day_parts(text: str) -> Optional[Tuple[int, int, int]]:
    raw = str(text or "").strip().lower()
    if not raw:
        return None

    m = re.search(r"(?:^|\b)(\d{1,2})(?::(\d{2}))?(?::(\d{2}))?\s*(am|pm)(?:\b|$)", raw)
    if m:
        hour = int(m.group(1))
        minute = int(m.group(2) or 0)
        second = int(m.group(3) or 0)
        meridiem = str(m.group(4) or "").lower()
        if hour < 1 or hour > 12 or minute > 59 or second > 59:
            return None
        if meridiem == "am":
            hour = 0 if hour == 12 else hour
        else:
            hour = 12 if hour == 12 else hour + 12
        return (hour, minute, second)

    m24 = re.search(r"(?:^|\b)([01]?\d|2[0-3]):([0-5]\d)(?::([0-5]\d))?(?:\b|$)", raw)
    if m24:
        return (int(m24.group(1)), int(m24.group(2)), int(m24.group(3) or 0))

    # Accept compact local times like "at 710" or "at 0710".
    m_compact = re.search(r"\bat\s+(\d{3,4})(?:\b|$)", raw)
    if m_compact:
        digits = str(m_compact.group(1) or "")
        if len(digits) == 3:
            hour = int(digits[0])
            minute = int(digits[1:])
        else:
            hour = int(digits[:2])
            minute = int(digits[2:])
        if 0 <= hour <= 23 and 0 <= minute <= 59:
            return (hour, minute, 0)

    # Accept common natural phrasing like "at 6" / "at 18" / "at 6:30"
    m_at = re.search(r"\bat\s+([01]?\d|2[0-3])(?::([0-5]\d))?(?::([0-5]\d))?(?:\b|$)", raw)
    if m_at:
        return (int(m_at.group(1)), int(m_at.group(2) or 0), int(m_at.group(3) or 0))

    # Accept "6 o'clock" style.
    m_oclock = re.search(r"\b([1-9]|1[0-2])\s*o'?clock\b", raw)
    if m_oclock:
        return (int(m_oclock.group(1)), 0, 0)

    return None


def _ai_tasks_infer_interval_from_text(text: str) -> float:
    raw = str(text or "").strip().lower()
    if not raw:
        return 0.0

    m = re.search(
        r"\bevery\s+(\d+)\s*(second|seconds|minute|minutes|hour|hours|day|days|week|weeks)\b",
        raw,
    )
    if m:
        count = int(m.group(1))
        unit = str(m.group(2) or "").lower()
        if count <= 0:
            return 0.0
        if unit.startswith("second"):
            return float(count)
        if unit.startswith("minute"):
            return float(count * 60)
        if unit.startswith("hour"):
            return float(count * 3600)
        if unit.startswith("day"):
            return float(count * 86400)
        if unit.startswith("week"):
            return float(count * 604800)

    if any(marker in raw for marker in AI_TASKS_DAILY_MARKERS):
        return 86400.0
    if "weekdays" in raw or "weekday" in raw:
        return 86400.0
    if any(marker in raw for marker in AI_TASKS_WEEKLY_MARKERS):
        return 604800.0
    return 0.0


def _ai_tasks_extract_weekdays(text: str) -> List[int]:
    raw = str(text or "").strip().lower()
    if not raw:
        return []
    if "weekdays" in raw or "weekday" in raw:
        return [0, 1, 2, 3, 4]
    if "weekends" in raw or "weekend" in raw:
        return [5, 6]

    out: List[int] = []
    for token, idx in AI_TASKS_WEEKDAY_MAP.items():
        if re.search(rf"\b{re.escape(token)}\b", raw):
            out.append(int(idx))
    if not out:
        return []
    return sorted(set(out))


def _ai_tasks_derive_recurrence(
    *,
    when_txt: Any,
    interval: float,
    next_run_ts: float,
    fallback_text: str = "",
) -> Dict[str, Any]:
    if interval <= 0 or next_run_ts <= 0:
        return {}

    text = str(when_txt or "").strip().lower()
    if not text and fallback_text:
        text = str(fallback_text).strip().lower()

    anchor_local = datetime.fromtimestamp(float(next_run_ts)).astimezone()
    weekdays = _ai_tasks_extract_weekdays(text)
    time_parts = _ai_tasks_extract_time_of_day_parts(text)

    if any(marker in text for marker in AI_TASKS_DAILY_MARKERS) or "weekday" in text or "weekdays" in text:
        recurrence: Dict[str, Any] = {
            "kind": "daily_local_time",
            "hour": int(anchor_local.hour),
            "minute": int(anchor_local.minute),
            "second": int(anchor_local.second),
        }
        if weekdays:
            recurrence["weekdays"] = weekdays
        return recurrence

    if any(marker in text for marker in AI_TASKS_WEEKLY_MARKERS) or weekdays:
        return {
            "kind": "weekly_local_time",
            "hour": int(anchor_local.hour),
            "minute": int(anchor_local.minute),
            "second": int(anchor_local.second),
            "weekdays": weekdays or [int(anchor_local.weekday())],
        }

    if time_parts and abs(float(interval) - 86400.0) <= 1.0:
        return {
            "kind": "daily_local_time",
            "hour": int(anchor_local.hour),
            "minute": int(anchor_local.minute),
            "second": int(anchor_local.second),
        }
    if time_parts and abs(float(interval) - 604800.0) <= 1.0:
        return {
            "kind": "weekly_local_time",
            "hour": int(anchor_local.hour),
            "minute": int(anchor_local.minute),
            "second": int(anchor_local.second),
            "weekdays": [int(anchor_local.weekday())],
        }
    return {}


def _ai_tasks_parse_when(when_ts: Any, when_txt: Any, in_seconds: Any) -> Optional[float]:
    now = time.time()

    if when_ts is not None:
        try:
            return float(when_ts)
        except Exception:
            pass

    if in_seconds is not None:
        try:
            return now + float(in_seconds)
        except Exception:
            pass

    if isinstance(when_txt, str) and when_txt.strip():
        text = when_txt.strip()
        text_lower = text.lower()
        if text.isdigit():
            # Treat long digit strings as epoch timestamps; short values are likely local hour-of-day.
            if len(text) >= 9:
                try:
                    return float(text)
                except Exception:
                    return None
            try:
                hour = int(text)
            except Exception:
                return None
            if 0 <= hour <= 23:
                now_dt = datetime.now().astimezone()
                dt = now_dt.replace(hour=hour, minute=0, second=0, microsecond=0)
                if dt.timestamp() <= now_dt.timestamp():
                    dt = dt + timedelta(days=1)
                return dt.timestamp()
            return None

        try:
            iso_text = text
            if iso_text.endswith("Z"):
                iso_text = iso_text[:-1] + "+00:00"
            dt = datetime.fromisoformat(iso_text)
        except Exception:
            dt = None

        if dt is None:
            for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M"):
                try:
                    dt = datetime.strptime(text, fmt)
                    break
                except Exception:
                    dt = None

        if dt is None:
            parts = _ai_tasks_extract_time_of_day_parts(text_lower)
            if parts:
                hour, minute, second = parts
                now_dt = datetime.now().astimezone()
                dt = now_dt.replace(hour=hour, minute=minute, second=second, microsecond=0)
                if dt.timestamp() <= now_dt.timestamp():
                    dt = dt + timedelta(days=1)

        if dt is None:
            return None
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=datetime.now().astimezone().tzinfo)
        return dt.timestamp()

    return None


def _ai_tasks_clean_task_prompt(task_text: Any) -> str:
    raw = str(task_text or "").strip()
    if not raw:
        return ""
    compact = re.sub(r"\s+", " ", raw).strip()
    compact = AI_TASKS_LOCAL_TZ_HINT_RE.sub("", compact).strip(" ,.-")
    compact = AI_TASKS_WEATHER_DEFAULT_HINT_RE.sub("", compact).strip(" ,.-")
    if not compact:
        return ""

    cleaned = compact
    for pattern in AI_TASKS_SCHEDULE_PREFIX_PATTERNS:
        match = pattern.match(cleaned)
        if not match:
            continue
        remainder = cleaned[match.end() :].strip(" ,.-")
        if remainder:
            cleaned = remainder
            break
    time_prefix = AI_TASKS_TIME_PREFIX_RE.match(cleaned)
    if time_prefix:
        remainder = cleaned[time_prefix.end() :].strip(" ,.-")
        if remainder:
            cleaned = remainder
    cleaned = re.sub(r"^(?:to\s+)", "", cleaned, flags=re.IGNORECASE).strip()
    return cleaned or compact


def _ai_tasks_default_title(task_prompt: str, recurrence: Dict[str, Any], interval: float) -> str:
    base = _ai_tasks_clean_task_prompt(task_prompt)
    if base:
        base = re.sub(
            r"^(?:please\s+)?(?:hey\s+\w+\s*,?\s*)?(?:can you|could you|would you|will you|do)\s+",
            "",
            base,
            flags=re.IGNORECASE,
        ).strip()
        base = re.sub(r"\s+", " ", base).strip(" .")
    recurrence_kind = str((recurrence or {}).get("kind") or "").strip().lower()
    cadence_prefix = ""
    if recurrence_kind == "daily_local_time":
        cadence_prefix = "Daily"
    elif recurrence_kind == "weekly_local_time":
        cadence_prefix = "Weekly"
    elif recurrence_kind == "monthly_local_time":
        cadence_prefix = "Monthly"
    elif float(interval or 0.0) > 0:
        cadence_prefix = "Recurring"

    lowered = base.lower()
    if any(token in lowered for token in ("weather", "forecast", "rain chance", "rain chances")):
        return f"{cadence_prefix} Weather Forecast".strip()

    if base:
        if len(base) > 80:
            base = base[:77].rstrip() + "..."
        return base[0].upper() + base[1:] if len(base) > 1 else base.upper()

    if cadence_prefix:
        return f"{cadence_prefix} AI task"
    return "Scheduled AI task"

class _WebpageInspectorParser(HTMLParser):
    def __init__(self, *, base_url: str, max_links: int, max_images: int):
        super().__init__(convert_charrefs=True)
        self.base_url = base_url
        self.max_links = max(1, int(max_links))
        self.max_images = max(1, int(max_images))
        self.title = ""
        self.description = ""
        self._in_title = False
        self._suppress_depth = 0
        self._anchor_href = ""
        self._anchor_text_parts: List[str] = []
        self.links: List[Dict[str, Any]] = []
        self.images: List[Dict[str, Any]] = []
        self._text_parts: List[str] = []

    def _attrs(self, attrs) -> Dict[str, str]:
        out: Dict[str, str] = {}
        for k, v in attrs or []:
            key = _as_text(k).strip().lower()
            if not key:
                continue
            out[key] = _as_text(v).strip()
        return out

    def _abs(self, ref: str) -> str:
        raw = _as_text(ref).strip()
        if not raw:
            return ""
        return urllib.parse.urljoin(self.base_url, raw)

    def _push_link(self, href: str, text: str, *, source: str = "a") -> None:
        if len(self.links) >= self.max_links:
            return
        url = self._abs(href)
        if not url:
            return
        self.links.append(
            {
                "url": url,
                "text": " ".join(_as_text(text).split()).strip(),
                "source": source,
            }
        )

    def _push_image(
        self,
        src: str,
        *,
        alt: str = "",
        title: str = "",
        class_name: str = "",
        id_name: str = "",
        source: str = "img",
    ) -> None:
        if len(self.images) >= self.max_images:
            return
        url = self._abs(src)
        if not url:
            return
        self.images.append(
            {
                "url": url,
                "alt": " ".join(_as_text(alt).split()).strip(),
                "title": " ".join(_as_text(title).split()).strip(),
                "class": " ".join(_as_text(class_name).split()).strip(),
                "id": " ".join(_as_text(id_name).split()).strip(),
                "source": source,
            }
        )

    def handle_starttag(self, tag: str, attrs) -> None:
        t = _as_text(tag).strip().lower()
        amap = self._attrs(attrs)
        if t in {"script", "style", "noscript"}:
            self._suppress_depth += 1
            return
        if t == "title":
            self._in_title = True
            return
        if t == "meta":
            meta_key = _as_text(amap.get("name") or amap.get("property")).strip().lower()
            content = _as_text(amap.get("content")).strip()
            if meta_key == "description" and content and not self.description:
                self.description = content
            if meta_key in {"og:image", "twitter:image", "twitter:image:src"} and content:
                self._push_image(content, source="meta")
            return
        if t == "link":
            rel = _as_text(amap.get("rel")).strip().lower()
            href = _as_text(amap.get("href")).strip()
            if not href:
                return
            if "icon" in rel or "apple-touch-icon" in rel:
                self._push_image(href, source="icon")
            elif rel in {"canonical", "alternate"}:
                self._push_link(href, "", source="link")
            return
        if t == "a":
            self._anchor_href = _as_text(amap.get("href")).strip()
            self._anchor_text_parts = []
            return
        if t == "img":
            self._push_image(
                _as_text(amap.get("src")).strip(),
                alt=amap.get("alt", ""),
                title=amap.get("title", ""),
                class_name=amap.get("class", ""),
                id_name=amap.get("id", ""),
                source="img",
            )

    def handle_endtag(self, tag: str) -> None:
        t = _as_text(tag).strip().lower()
        if t in {"script", "style", "noscript"}:
            self._suppress_depth = max(0, self._suppress_depth - 1)
            return
        if t == "title":
            self._in_title = False
            return
        if t == "a":
            if self._anchor_href:
                text = " ".join(" ".join(self._anchor_text_parts).split()).strip()
                self._push_link(self._anchor_href, text, source="a")
            self._anchor_href = ""
            self._anchor_text_parts = []

    def handle_data(self, data: str) -> None:
        text = " ".join(_as_text(data).split()).strip()
        if not text:
            return
        if self._in_title and not self.title:
            self.title = text
        if self._suppress_depth > 0:
            return
        if self._anchor_href:
            self._anchor_text_parts.append(text)
        if len(self._text_parts) < 400:
            self._text_parts.append(text)

    def visible_text(self, *, max_chars: int = 1200) -> str:
        joined = " ".join(self._text_parts).strip()
        if len(joined) <= max_chars:
            return joined
        clipped = joined[:max_chars]
        if " " in clipped[200:]:
            clipped = clipped[: clipped.rfind(" ")]
        return clipped.rstrip(" .,;:") + "..."


def _score_image_candidate(image: Dict[str, Any]) -> int:
    words = " ".join(
        [
            _as_text(image.get("alt")).lower(),
            _as_text(image.get("title")).lower(),
            _as_text(image.get("class")).lower(),
            _as_text(image.get("id")).lower(),
            _as_text(image.get("url")).lower(),
            _as_text(image.get("source")).lower(),
        ]
    )
    score = 0
    if any(token in words for token in ("logo", "brand", "wordmark", "logotype")):
        score += 4
    if any(token in words for token in ("icon", "favicon")):
        score += 2
    if ".svg" in words:
        score += 1
    if any(token in words for token in ("sprite", "blank", "placeholder", "spacer")):
        score -= 3
    return score


def inspect_webpage(
    url: str,
    *,
    max_bytes: Optional[int] = None,
    timeout_sec: int = 20,
    max_links: int = 20,
    max_images: int = 20,
    platform: Optional[str] = None,
    origin: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    del platform, origin
    normalized_url = _normalize_url_input(url)
    err = _validate_url(normalized_url)
    if err:
        return {"tool": "inspect_webpage", "ok": False, "error": err}
    fetch_payload = _fetch_webpage(
        normalized_url,
        timeout_sec=int(timeout_sec or WEBPAGE_TIMEOUT_SEC),
        max_bytes=_coerce_int(
            max_bytes if max_bytes is not None else WEBPAGE_MAX_RESPONSE_BYTES,
            default=WEBPAGE_MAX_RESPONSE_BYTES,
            min_value=65_536,
            max_value=25_000_000,
        ),
    )

    fetch_attempts = fetch_payload.get("attempts") if isinstance(fetch_payload, dict) else []
    attempts_out: List[Dict[str, Any]] = []
    for row in fetch_attempts if isinstance(fetch_attempts, list) else []:
        if not isinstance(row, dict):
            continue
        attempts_out.append(
            {
                "status_code": int(row.get("status_code") or 0),
                "blocked": bool(row.get("blocked")),
                "error": str(row.get("error") or "").strip(),
            }
        )

    if not bool(fetch_payload.get("ok")):
        return {
            "tool": "inspect_webpage",
            "ok": False,
            "url": str(fetch_payload.get("url") or normalized_url),
            "error": str(fetch_payload.get("error") or "Unable to fetch webpage."),
            "fetch_attempts": attempts_out,
        }

    final_url = str(fetch_payload.get("url") or normalized_url)
    content_type = str(fetch_payload.get("content_type") or "")
    status_code = int(fetch_payload.get("status_code") or 0)
    raw_text = str(fetch_payload.get("raw_text") or "")
    truncated = bool(fetch_payload.get("truncated"))
    byte_count = int(fetch_payload.get("bytes") or len(raw_text.encode("utf-8", errors="ignore")))

    content_type_norm = content_type.split(";", 1)[0].strip().lower()
    is_html = (
        content_type_norm.startswith("text/html")
        or content_type_norm.startswith("application/xhtml+xml")
        or ("<html" in raw_text[:1200].lower())
    )

    parser = _WebpageInspectorParser(
        base_url=final_url,
        max_links=max_links,
        max_images=max_images,
    )
    if is_html:
        try:
            parser.feed(raw_text)
        except Exception:
            pass
        try:
            parser.close()
        except Exception:
            pass

    title_text = _as_text(parser.title).strip()
    description_text = _as_text(parser.description).strip()
    if is_html:
        content_text = _webpage_text_from_html(raw_text, max_chars=WEBPAGE_MAX_TEXT_CHARS)
    else:
        content_text = " ".join(raw_text.split()).strip()
    if not content_text and is_html:
        content_text = parser.visible_text(max_chars=WEBPAGE_MAX_TEXT_CHARS)

    text_preview = _webpage_preview(content_text, max_chars=WEBPAGE_MAX_PREVIEW_CHARS)
    bot_blocked = _webpage_looks_bot_blocked(
        status_code=status_code,
        title=title_text,
        description=description_text,
        preview=text_preview,
    )

    if bot_blocked and len(content_text) < 220:
        return {
            "tool": "inspect_webpage",
            "ok": False,
            "url": final_url,
            "status_code": status_code,
            "content_type": content_type,
            "bytes": byte_count,
            "truncated": truncated,
            "blocked_by_bot": True,
            "error": (
                "Webpage appears to be protected by anti-bot checks and no readable content was extracted. "
                "Try another source URL."
            ),
            "fetch_attempts": attempts_out,
        }

    unique_links: List[Dict[str, Any]] = []
    seen_links = set()
    for item in parser.links:
        link_url = _as_text(item.get("url")).strip()
        if not link_url or link_url in seen_links:
            continue
        seen_links.add(link_url)
        unique_links.append(item)

    unique_images: List[Dict[str, Any]] = []
    seen_images = set()
    for item in parser.images:
        image_url = _as_text(item.get("url")).strip()
        if not image_url or image_url in seen_images:
            continue
        seen_images.add(image_url)
        scored = dict(item)
        score = _score_image_candidate(scored)
        scored["score"] = score
        scored["logo_hint"] = bool(score >= 2)
        unique_images.append(scored)

    best_image_url = ""
    if unique_images:
        ranked = sorted(
            enumerate(unique_images),
            key=lambda pair: (int(pair[1].get("score") or 0), -pair[0]),
            reverse=True,
        )
        best_image_url = _as_text(ranked[0][1].get("url")).strip()

    content_chunk, content_meta = _slice_content(content_text, start=0, max_chars=WEBPAGE_MAX_TEXT_CHARS)

    return {
        "tool": "inspect_webpage",
        "ok": True,
        "url": final_url,
        "status_code": status_code,
        "content_type": content_type,
        "bytes": byte_count,
        "truncated": truncated,
        "blocked_by_bot": bool(bot_blocked),
        "title": title_text,
        "description": description_text,
        "text_preview": text_preview,
        "content": content_chunk,
        "content_meta": content_meta,
        "links": unique_links,
        "link_count": len(unique_links),
        "images": unique_images,
        "image_count": len(unique_images),
        "best_image_url": best_image_url or None,
        "fetch_attempts": attempts_out,
    }


def download_file(
    url: str,
    *,
    filename: Optional[str] = None,
    subdir: Optional[str] = None,
    max_bytes: Optional[int] = None,
    timeout_sec: int = 30,
    platform: Optional[str] = None,
    origin: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    _ensure_dirs()
    normalized_url = _normalize_url_input(url)
    err = _validate_url(normalized_url)
    if err:
        return {"tool": "download_file", "ok": False, "error": err}

    # Resolve target directory inside workspace root (default: downloads)
    target_dir = AGENT_LAB_DIR / (subdir or "downloads")
    try:
        target_dir = target_dir.resolve()
    except Exception:
        target_dir = AGENT_DOCUMENTS_DIR
    if not (
        target_dir == AGENT_LAB_DIR.resolve()
        or AGENT_LAB_DIR.resolve() in target_dir.parents
    ):
        return {"tool": "download_file", "ok": False, "error": "Target directory not allowed."}
    target_dir.mkdir(parents=True, exist_ok=True)

    parsed = urllib.parse.urlparse(normalized_url)
    default_name = _sanitize_filename(os.path.basename(parsed.path)) or "download.bin"
    safe_name = _sanitize_filename(filename or default_name) or "download.bin"
    dest = target_dir / safe_name

    import hashlib

    hasher = hashlib.sha256()
    size = 0
    content_type = ""
    try:
        req = urllib.request.Request(normalized_url, headers={"User-Agent": "Tater/1.0"})
        with urllib.request.urlopen(req, timeout=timeout_sec) as resp:
            content_type = resp.headers.get("Content-Type", "") or ""
            final_url = resp.geturl() or normalized_url
            with dest.open("wb") as f:
                while True:
                    chunk = resp.read(8192)
                    if not chunk:
                        break
                    size += len(chunk)
                    f.write(chunk)
                    hasher.update(chunk)
    except Exception as e:
        try:
            if dest.exists():
                dest.unlink()
        except Exception:
            pass
        return {"tool": "download_file", "ok": False, "error": str(e)}

    _log_write("download_file", dest, size)
    out: Dict[str, Any] = {
        "tool": "download_file",
        "ok": True,
        "url": final_url,
        "path": _display_workspace_path(dest),
        "bytes": size,
        "sha256": hasher.hexdigest(),
        "content_type": content_type,
    }

    media_type, detected_mime = _download_file_detect_media(dest, content_type)
    out["artifact"] = {
        "type": media_type,
        "path": _display_workspace_path(dest),
        "name": dest.name,
        "mimetype": detected_mime,
        "source": "download_file",
        "size": size,
    }

    return out


def read_file(path: str, start: int = 0, max_chars: Optional[int] = None) -> Dict[str, Any]:
    _ensure_dirs()
    allowed = [AGENT_LAB_DIR]
    resolved = _resolve_safe_path(path, allowed)
    if not resolved:
        return {"tool": "read_file", "ok": False, "error": "Path not allowed."}
    if not resolved.exists() or not resolved.is_file():
        return {"tool": "read_file", "ok": False, "error": "File not found."}
    try:
        full_content, metadata = _extract_file_content(resolved)
        chunk, window = _slice_content(full_content, start=start, max_chars=max_chars)
        source_truncated = bool(metadata.get("source_truncated"))
        return {
            "tool": "read_file",
            "ok": True,
            "path": _display_workspace_path(resolved),
            "content": chunk,
            "source_truncated": source_truncated,
            "truncated": bool(source_truncated or window.get("has_more")),
            **window,
            **metadata,
        }
    except Exception as e:
        return {"tool": "read_file", "ok": False, "error": str(e)}


def _is_hidden_path(path: Path) -> bool:
    for part in path.parts:
        if part.startswith(".") and part not in {".", ".."}:
            return True
    return False


def _find_query_hits(
    content: str,
    query: str,
    *,
    case_sensitive: bool = False,
    max_hits: int = 50,
) -> List[Dict[str, Any]]:
    if not query:
        return []
    needle = query if case_sensitive else query.lower()
    hits: List[Dict[str, Any]] = []
    for idx, line in enumerate(content.splitlines(), start=1):
        hay = line if case_sensitive else line.lower()
        if needle in hay:
            snippet = line.strip()
            if len(snippet) > 320:
                snippet = snippet[:320].rstrip() + "..."
            hits.append({"line": idx, "snippet": snippet})
            if len(hits) >= max_hits:
                break
    return hits


def search_files(
    query: str,
    *,
    path: Optional[str] = None,
    max_results: int = _SEARCH_DEFAULT_MAX_RESULTS,
    case_sensitive: bool = False,
    include_hidden: bool = False,
    file_glob: Optional[str] = None,
) -> Dict[str, Any]:
    _ensure_dirs()
    needle = str(query or "").strip()
    if not needle:
        return {"tool": "search_files", "ok": False, "error": "query is required."}

    max_results_i = _coerce_int(max_results, default=_SEARCH_DEFAULT_MAX_RESULTS, min_value=1, max_value=2000)
    allow_roots = [AGENT_LAB_DIR]

    try:
        targets: List[Path] = []
        if path and str(path).strip():
            resolved = _resolve_safe_path(str(path), allow_roots)
            if not resolved:
                return {"tool": "search_files", "ok": False, "error": "Path not allowed."}
            if not resolved.exists():
                return {"tool": "search_files", "ok": False, "error": "Path not found."}
            targets = [resolved]
        else:
            targets = [AGENT_DOCUMENTS_DIR, AGENT_DOWNLOADS_DIR, AGENT_WORKSPACE_DIR]

        scanned_files = 0
        skipped_files = 0
        matched_files = 0
        results: List[Dict[str, Any]] = []
        seen: set[str] = set()
        per_file_hit_limit = 5

        def _iter_files(base: Path) -> List[Path]:
            if base.is_file():
                return [base]
            out: List[Path] = []
            for p in base.rglob("*"):
                if p.is_file():
                    out.append(p)
            return out

        for target in targets:
            files = _iter_files(target)
            for file_path in files:
                file_key = str(file_path.resolve())
                if file_key in seen:
                    continue
                seen.add(file_key)
                if not include_hidden and _is_hidden_path(file_path):
                    continue
                if file_glob and not fnmatch.fnmatch(file_path.name, str(file_glob)):
                    continue

                scanned_files += 1
                try:
                    text, _meta = _extract_file_content(file_path)
                except Exception:
                    skipped_files += 1
                    continue
                if len(text) > _SEARCH_MAX_FILE_CHARS:
                    text = text[:_SEARCH_MAX_FILE_CHARS]

                hits = _find_query_hits(
                    text,
                    needle,
                    case_sensitive=case_sensitive,
                    max_hits=per_file_hit_limit,
                )
                if not hits:
                    continue

                matched_files += 1
                for hit in hits:
                    results.append(
                        {
                            "path": _display_workspace_path(file_path),
                            "line": hit["line"],
                            "snippet": hit["snippet"],
                        }
                    )
                    if len(results) >= max_results_i:
                        break
                if len(results) >= max_results_i:
                    break
            if len(results) >= max_results_i:
                break

        return {
            "tool": "search_files",
            "ok": True,
            "query": needle,
            "results": results,
            "count": len(results),
            "scanned_files": scanned_files,
            "matched_files": matched_files,
            "skipped_files": skipped_files,
            "paths": [_display_workspace_path(p) for p in targets],
            "max_results": max_results_i,
        }
    except Exception as e:
        return {"tool": "search_files", "ok": False, "error": str(e)}


def write_file(
    path: str,
    content: Optional[str] = None,
    *,
    content_b64: Optional[str] = None,
    content_lines: Optional[List[str]] = None,
) -> Dict[str, Any]:
    _ensure_dirs()
    allowed = [AGENT_LAB_DIR]
    resolved = _resolve_safe_path(path, allowed)
    if not resolved:
        return {"tool": "write_file", "ok": False, "error": "Path not allowed."}
    try:
        # Prevent creating executable modules via write_file.
        if resolved.suffix == ".py":
            try:
                rel = resolved.relative_to(AGENT_LAB_DIR)
                top = rel.parts[0] if rel.parts else ""
            except Exception:
                top = ""
            if top in {"verba", "portals"}:
                return {
                    "tool": "write_file",
                    "ok": False,
                    "error": "Direct python writes are disabled for verba/portals.",
                }
            return {
                "tool": "write_file",
                "ok": False,
                "error": "Python files are not allowed via write_file.",
            }
        resolved.parent.mkdir(parents=True, exist_ok=True)
        if content_b64:
            try:
                import base64
                data = base64.b64decode(content_b64.encode("utf-8")).decode("utf-8")
            except Exception as e:
                return {"tool": "write_file", "ok": False, "error": f"Invalid content_b64: {e}"}
        elif isinstance(content_lines, list):
            data = "\n".join(str(x) for x in content_lines)
        else:
            data = content if content is not None else ""
        resolved.write_text(data, encoding="utf-8")
        _log_write("write_file", resolved, len(data.encode("utf-8")))
        return {"tool": "write_file", "ok": True, "path": _display_workspace_path(resolved), "bytes": len(data)}
    except Exception as e:
        return {"tool": "write_file", "ok": False, "error": str(e)}


def list_directory(path: str) -> Dict[str, Any]:
    _ensure_dirs()
    allowed = [AGENT_LAB_DIR]
    resolved = _resolve_safe_path(path, allowed)
    if not resolved:
        return {"tool": "list_directory", "ok": False, "error": "Path not allowed."}
    if not resolved.exists() or not resolved.is_dir():
        return {"tool": "list_directory", "ok": False, "error": "Directory not found."}
    try:
        files = []
        dirs = []
        for item in sorted(resolved.iterdir()):
            if item.is_dir():
                dirs.append(item.name)
            else:
                files.append(item.name)
        return {"tool": "list_directory", "ok": True, "path": _display_workspace_path(resolved), "files": files, "directories": dirs}
    except Exception as e:
        return {"tool": "list_directory", "ok": False, "error": str(e)}


def _safe_archive_target(base_dir: Path, member_name: str) -> Optional[Path]:
    raw = str(member_name or "").strip().replace("\\", "/")
    if not raw:
        return None
    while raw.startswith("/"):
        raw = raw[1:]
    normalized = os.path.normpath(raw).replace("\\", "/")
    if normalized in {"", ".", ".."} or normalized.startswith("../"):
        return None
    first = normalized.split("/", 1)[0]
    if ":" in first:
        return None
    try:
        target = (base_dir / normalized).resolve()
        base = base_dir.resolve()
    except Exception:
        return None
    if target == base or base in target.parents:
        return target
    return None


def list_archive(path: str, max_entries: int = _ARCHIVE_LIST_MAX_ENTRIES) -> Dict[str, Any]:
    _ensure_dirs()
    allowed = [AGENT_LAB_DIR]
    archive_path = _resolve_safe_path(path, allowed)
    if not archive_path:
        return {"tool": "list_archive", "ok": False, "error": "Path not allowed."}
    if not archive_path.exists() or not archive_path.is_file():
        return {"tool": "list_archive", "ok": False, "error": "Archive not found."}

    max_entries_i = _coerce_int(max_entries, default=_ARCHIVE_LIST_MAX_ENTRIES, min_value=1, max_value=5000)
    entries: List[Dict[str, Any]] = []
    truncated = False
    lowered_name = archive_path.name.lower()

    try:
        if lowered_name.endswith(".7z"):
            try:
                import py7zr
            except Exception as e:
                return {"tool": "list_archive", "ok": False, "error": "7z support requires the `py7zr` package."}
            with py7zr.SevenZipFile(str(archive_path), mode="r") as zf:
                for info in zf.list():
                    if len(entries) >= max_entries_i:
                        truncated = True
                        break
                    name = str(getattr(info, "filename", "") or "")
                    is_dir = bool(getattr(info, "is_directory", False)) or name.endswith("/")
                    entries.append(
                        {
                            "name": name,
                            "size": int(getattr(info, "uncompressed", 0) or 0),
                            "compressed_size": int(getattr(info, "compressed", 0) or 0),
                            "is_dir": is_dir,
                        }
                    )
            return {
                "tool": "list_archive",
                "ok": True,
                "path": _display_workspace_path(archive_path),
                "format": "7z",
                "entries": entries,
                "count": len(entries),
                "truncated": truncated,
            }

        if lowered_name.endswith(".rar"):
            try:
                import rarfile
            except Exception as e:
                return {"tool": "list_archive", "ok": False, "error": "RAR support requires the `rarfile` package."}
            with rarfile.RarFile(str(archive_path), "r") as rf:
                for info in rf.infolist():
                    if len(entries) >= max_entries_i:
                        truncated = True
                        break
                    entries.append(
                        {
                            "name": info.filename,
                            "size": int(getattr(info, "file_size", 0) or 0),
                            "compressed_size": int(getattr(info, "compress_size", 0) or 0),
                            "is_dir": bool(info.isdir()),
                            "is_symlink": bool(getattr(info, "is_symlink", lambda: False)()),
                        }
                    )
            return {
                "tool": "list_archive",
                "ok": True,
                "path": _display_workspace_path(archive_path),
                "format": "rar",
                "entries": entries,
                "count": len(entries),
                "truncated": truncated,
            }

        if zipfile.is_zipfile(archive_path):
            with zipfile.ZipFile(archive_path, "r") as zf:
                for info in zf.infolist():
                    if len(entries) >= max_entries_i:
                        truncated = True
                        break
                    entries.append(
                        {
                            "name": info.filename,
                            "size": int(info.file_size),
                            "compressed_size": int(info.compress_size),
                            "is_dir": bool(info.is_dir()),
                        }
                    )
            return {
                "tool": "list_archive",
                "ok": True,
                "path": _display_workspace_path(archive_path),
                "format": "zip",
                "entries": entries,
                "count": len(entries),
                "truncated": truncated,
            }

        if tarfile.is_tarfile(archive_path):
            with tarfile.open(archive_path, "r:*") as tf:
                for member in tf.getmembers():
                    if len(entries) >= max_entries_i:
                        truncated = True
                        break
                    entries.append(
                        {
                            "name": member.name,
                            "size": int(member.size or 0),
                            "is_dir": bool(member.isdir()),
                            "is_symlink": bool(member.issym() or member.islnk()),
                        }
                    )
            return {
                "tool": "list_archive",
                "ok": True,
                "path": _display_workspace_path(archive_path),
                "format": "tar",
                "entries": entries,
                "count": len(entries),
                "truncated": truncated,
            }

        return {
            "tool": "list_archive",
            "ok": False,
            "error": "Unsupported archive format. Supported: zip, tar(.gz/.bz2/.xz), 7z, rar.",
        }
    except Exception as e:
        return {"tool": "list_archive", "ok": False, "error": str(e)}


def _archive_output_folder_name(path: Path) -> str:
    name = path.name
    lowered = name.lower()
    suffixes = [".tar.gz", ".tgz", ".tar.bz2", ".tbz2", ".tar.xz", ".txz", ".zip", ".tar", ".7z", ".rar"]
    for suffix in suffixes:
        if lowered.endswith(suffix):
            base = name[: -len(suffix)]
            return base or path.stem or "archive"
    return path.stem or "archive"


def extract_archive(
    path: str,
    *,
    destination: Optional[str] = None,
    overwrite: bool = False,
    max_files: int = _ARCHIVE_EXTRACT_MAX_FILES,
    max_total_bytes: Optional[int] = None,
) -> Dict[str, Any]:
    _ensure_dirs()
    allowed = [AGENT_LAB_DIR]
    archive_path = _resolve_safe_path(path, allowed)
    if not archive_path:
        return {"tool": "extract_archive", "ok": False, "error": "Path not allowed."}
    if not archive_path.exists() or not archive_path.is_file():
        return {"tool": "extract_archive", "ok": False, "error": "Archive not found."}

    default_dest = f"workspace/extracted_{_archive_output_folder_name(archive_path)}"
    auto_destination = destination in (None, "")
    dest_path = _resolve_safe_path(destination or default_dest, allowed)
    if not dest_path:
        return {"tool": "extract_archive", "ok": False, "error": "Destination path not allowed."}
    if dest_path.exists() and not dest_path.is_dir():
        return {"tool": "extract_archive", "ok": False, "error": "Destination exists and is not a directory."}
    dest_existed_before = dest_path.exists()
    if not dest_existed_before:
        dest_path.mkdir(parents=True, exist_ok=True)

    max_files_i = _coerce_int(max_files, default=_ARCHIVE_EXTRACT_MAX_FILES, min_value=1, max_value=20000)

    extracted: List[str] = []
    skipped: List[Dict[str, str]] = []
    extracted_count = 0
    total_bytes = 0
    limit_hit = False
    lowered_name = archive_path.name.lower()

    def _record_skip(name: str, reason: str) -> None:
        skipped.append({"name": name, "reason": reason})

    try:
        if lowered_name.endswith(".7z"):
            try:
                import py7zr
            except Exception as e:
                return {"tool": "extract_archive", "ok": False, "error": "7z support requires the `py7zr` package."}
            with py7zr.SevenZipFile(str(archive_path), mode="r") as zf:
                selected_names: List[str] = []
                selected_targets: List[Path] = []
                selected_sizes: List[int] = []
                for info in zf.list():
                    name = str(getattr(info, "filename", "") or "")
                    is_dir = bool(getattr(info, "is_directory", False)) or name.endswith("/")
                    if is_dir:
                        continue
                    target = _safe_archive_target(dest_path, name)
                    if not target:
                        _record_skip(name, "unsafe_path")
                        continue
                    size = int(getattr(info, "uncompressed", 0) or 0)
                    if size < 0:
                        _record_skip(name, "invalid_size")
                        continue
                    if extracted_count >= max_files_i:
                        _record_skip(name, "max_files_reached")
                        limit_hit = True
                        break
                    if target.exists() and not overwrite:
                        _record_skip(name, "exists")
                        continue
                    selected_names.append(name)
                    selected_targets.append(target)
                    selected_sizes.append(size)
                    extracted_count += 1
                    total_bytes += size
                    extracted.append(_display_workspace_path(target))
                if selected_names:
                    zf.extract(path=str(dest_path), targets=selected_names)
                    for target, size in zip(selected_targets, selected_sizes):
                        _log_write("extract_archive", target, size)

            return {
                "tool": "extract_archive",
                "ok": True,
                "path": _display_workspace_path(archive_path),
                "format": "7z",
                "destination": _display_workspace_path(dest_path),
                "extracted_count": extracted_count,
                "extracted": extracted,
                "skipped_count": len(skipped),
                "skipped": skipped,
                "bytes_written": total_bytes,
                "limit_hit": limit_hit,
            }

        if lowered_name.endswith(".rar"):
            try:
                import rarfile
            except Exception as e:
                return {"tool": "extract_archive", "ok": False, "error": "RAR support requires the `rarfile` package."}
            with rarfile.RarFile(str(archive_path), "r") as rf:
                for info in rf.infolist():
                    name = info.filename
                    if info.isdir():
                        continue
                    is_symlink = bool(getattr(info, "is_symlink", lambda: False)())
                    if is_symlink:
                        _record_skip(name, "symlink_not_allowed")
                        continue
                    target = _safe_archive_target(dest_path, name)
                    if not target:
                        _record_skip(name, "unsafe_path")
                        continue
                    size = int(getattr(info, "file_size", 0) or 0)
                    if size < 0:
                        _record_skip(name, "invalid_size")
                        continue
                    if extracted_count >= max_files_i:
                        _record_skip(name, "max_files_reached")
                        limit_hit = True
                        break
                    if target.exists() and not overwrite:
                        _record_skip(name, "exists")
                        continue
                    target.parent.mkdir(parents=True, exist_ok=True)
                    with rf.open(info) as src, target.open("wb") as dst:
                        shutil.copyfileobj(src, dst)
                    total_bytes += size
                    extracted_count += 1
                    extracted.append(_display_workspace_path(target))
                    _log_write("extract_archive", target, size)

            return {
                "tool": "extract_archive",
                "ok": True,
                "path": _display_workspace_path(archive_path),
                "format": "rar",
                "destination": _display_workspace_path(dest_path),
                "extracted_count": extracted_count,
                "extracted": extracted,
                "skipped_count": len(skipped),
                "skipped": skipped,
                "bytes_written": total_bytes,
                "limit_hit": limit_hit,
            }

        if zipfile.is_zipfile(archive_path):
            with zipfile.ZipFile(archive_path, "r") as zf:
                for info in zf.infolist():
                    if info.is_dir():
                        continue
                    name = info.filename
                    target = _safe_archive_target(dest_path, name)
                    if not target:
                        _record_skip(name, "unsafe_path")
                        continue
                    if info.file_size < 0:
                        _record_skip(name, "invalid_size")
                        continue
                    if extracted_count >= max_files_i:
                        _record_skip(name, "max_files_reached")
                        limit_hit = True
                        break
                    if target.exists() and not overwrite:
                        _record_skip(name, "exists")
                        continue
                    target.parent.mkdir(parents=True, exist_ok=True)
                    with zf.open(info, "r") as src, target.open("wb") as dst:
                        shutil.copyfileobj(src, dst)
                    total_bytes += int(info.file_size)
                    extracted_count += 1
                    extracted.append(_display_workspace_path(target))
                    _log_write("extract_archive", target, int(info.file_size))

            return {
                "tool": "extract_archive",
                "ok": True,
                "path": _display_workspace_path(archive_path),
                "format": "zip",
                "destination": _display_workspace_path(dest_path),
                "extracted_count": extracted_count,
                "extracted": extracted,
                "skipped_count": len(skipped),
                "skipped": skipped,
                "bytes_written": total_bytes,
                "limit_hit": limit_hit,
            }

        if tarfile.is_tarfile(archive_path):
            with tarfile.open(archive_path, "r:*") as tf:
                for member in tf.getmembers():
                    name = member.name
                    if member.isdir():
                        continue
                    if member.issym() or member.islnk():
                        _record_skip(name, "symlink_not_allowed")
                        continue
                    target = _safe_archive_target(dest_path, name)
                    if not target:
                        _record_skip(name, "unsafe_path")
                        continue
                    size = int(member.size or 0)
                    if extracted_count >= max_files_i:
                        _record_skip(name, "max_files_reached")
                        limit_hit = True
                        break
                    if target.exists() and not overwrite:
                        _record_skip(name, "exists")
                        continue
                    stream = tf.extractfile(member)
                    if stream is None:
                        _record_skip(name, "unreadable_entry")
                        continue
                    target.parent.mkdir(parents=True, exist_ok=True)
                    with stream as src, target.open("wb") as dst:
                        shutil.copyfileobj(src, dst)
                    total_bytes += size
                    extracted_count += 1
                    extracted.append(_display_workspace_path(target))
                    _log_write("extract_archive", target, size)

            return {
                "tool": "extract_archive",
                "ok": True,
                "path": _display_workspace_path(archive_path),
                "format": "tar",
                "destination": _display_workspace_path(dest_path),
                "extracted_count": extracted_count,
                "extracted": extracted,
                "skipped_count": len(skipped),
                "skipped": skipped,
                "bytes_written": total_bytes,
                "limit_hit": limit_hit,
            }

        return {
            "tool": "extract_archive",
            "ok": False,
            "error": "Unsupported archive format. Supported: zip, tar(.gz/.bz2/.xz), 7z, rar.",
        }
    except Exception as e:
        return {"tool": "extract_archive", "ok": False, "error": str(e)}
    finally:
        # Clean up auto-created default extraction dirs when nothing was written.
        if auto_destination and not dest_existed_before:
            try:
                if (
                    dest_path.exists()
                    and dest_path.is_dir()
                    and dest_path != AGENT_WORKSPACE_DIR
                    and next(dest_path.iterdir(), None) is None
                ):
                    dest_path.rmdir()
            except Exception:
                pass


def delete_file(path: str) -> Dict[str, Any]:
    _ensure_dirs()
    allowed = [AGENT_LAB_DIR]
    resolved = _resolve_safe_path(path, allowed)
    if not resolved:
        return {"tool": "delete_file", "ok": False, "error": "Path not allowed."}
    if not resolved.exists() or not resolved.is_file():
        return {"tool": "delete_file", "ok": False, "error": "File not found."}
    try:
        resolved.unlink()
        _log_write("delete_file", resolved, 0)
        return {"tool": "delete_file", "ok": True, "path": _display_workspace_path(resolved), "deleted": True}
    except Exception as e:
        return {"tool": "delete_file", "ok": False, "error": str(e)}

def _exp_plugin_path(name: str) -> Path:
    return AGENT_PLUGINS_DIR / f"{name}.py"


def promote_plugin(name: str, confirm: Optional[bool] = None, delete_source: bool = False) -> Dict[str, Any]:
    _ensure_dirs()
    if not confirm:
        return {
            "tool": "promote_plugin",
            "ok": False,
            "error": "Confirmation required.",
            "needs": ["Please confirm promotion to stable verbas by setting confirm=true."],
        }
    if not _SAFE_NAME_RE.fullmatch(name or ""):
        return {"tool": "promote_plugin", "ok": False, "error": "Invalid plugin name."}
    src = _exp_plugin_path(name)
    if not src.exists():
        return {"tool": "promote_plugin", "ok": False, "error": "Plugin not found."}
    dest = STABLE_PLUGINS_DIR / f"{name}.py"
    try:
        dest.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(src, dest)
        reload_verbas()
        _log_write("promote_plugin", dest, dest.stat().st_size if dest.exists() else 0)
        if delete_source:
            src.unlink()
        return {"tool": "promote_plugin", "ok": True, "path": _display_workspace_path(dest)}
    except Exception as e:
        return {"tool": "promote_plugin", "ok": False, "error": str(e)}

def _as_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (bytes, bytearray)):
        try:
            return value.decode("utf-8", errors="replace")
        except Exception:
            return str(value)
    return str(value)


def _origin_value(origin: Optional[Dict[str, Any]], *keys: str) -> str:
    if not isinstance(origin, dict):
        return ""
    for key in keys:
        val = origin.get(key)
        if val in (None, ""):
            continue
        text = _as_text(val).strip()
        if text:
            return text
    return ""


def _normalize_key_segment(value: Any, *, default: str) -> str:
    raw = _as_text(value).strip().lower()
    if not raw:
        raw = default
    cleaned = re.sub(r"[^a-z0-9_.:\-]+", "_", raw).strip("_")
    return cleaned or default



def write_workspace_note(content: str) -> Dict[str, Any]:
    _ensure_dirs()
    ts = time.strftime("%Y%m%d_%H%M%S")
    suffix = uuid.uuid4().hex[:8]
    filename = f"note_{ts}_{suffix}.md"
    path = AGENT_WORKSPACE_DIR / filename
    try:
        data = content or ""
        path.write_text(data, encoding="utf-8")
        _log_write("write_workspace_note", path, len(data.encode("utf-8")))
        return {"tool": "write_workspace_note", "ok": True, "path": _display_workspace_path(path)}
    except Exception as e:
        return {"tool": "write_workspace_note", "ok": False, "error": str(e)}


def list_workspace() -> Dict[str, Any]:
    _ensure_dirs()
    try:
        files = sorted([p.name for p in AGENT_WORKSPACE_DIR.iterdir() if p.is_file()])
        return {"tool": "list_workspace", "ok": True, "files": files}
    except Exception as e:
        return {"tool": "list_workspace", "ok": False, "error": str(e)}
